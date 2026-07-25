"""Text extraction tests for document file types."""

from pathlib import Path

import pytest

from engine import (
    DOCUMENT_EXTENSIONS,
    extract_text_csv,
    extract_text_docx,
    extract_text_from_file,
    extract_text_pdf,
    extract_text_plain,
    extract_text_pptx,
    extract_text_xlsx,
)


class TestExtractTextPlain:
    def test_extract_txt(self, tmp_path: Path) -> None:
        """Plain .txt files return their content."""
        f = tmp_path / "notes.txt"
        f.write_text("Hello world\nSecond line", encoding="utf-8")
        result = extract_text_plain(f)
        assert result == "Hello world\nSecond line"

    def test_extract_md(self, tmp_path: Path) -> None:
        """Markdown files return their content."""
        f = tmp_path / "readme.md"
        f.write_text("# Title\n\nBody text", encoding="utf-8")
        result = extract_text_plain(f)
        assert "# Title" in result
        assert "Body text" in result

    def test_empty_file_returns_none(self, tmp_path: Path) -> None:
        """Empty files return None."""
        f = tmp_path / "empty.txt"
        f.write_text("", encoding="utf-8")
        result = extract_text_plain(f)
        assert result is None

    def test_missing_file_returns_none(self, tmp_path: Path) -> None:
        """Missing files return None."""
        result = extract_text_plain(tmp_path / "nonexistent.txt")
        assert result is None


class TestExtractTextCSV:
    def test_extract_csv(self, tmp_path: Path) -> None:
        """CSV files return pipe-delimited rows."""
        f = tmp_path / "data.csv"
        f.write_text("name,age\nAlice,30\nBob,25", encoding="utf-8")
        result = extract_text_csv(f)
        assert result is not None
        assert "name | age" in result
        assert "Alice | 30" in result

    def test_empty_csv_returns_none(self, tmp_path: Path) -> None:
        """Empty CSV returns None."""
        f = tmp_path / "empty.csv"
        f.write_text("", encoding="utf-8")
        result = extract_text_csv(f)
        assert result is None


class TestExtractTextXLSX:
    def test_extract_xlsx(self, tmp_path: Path) -> None:
        """XLSX files return sheet content with headers."""
        pytest.importorskip("openpyxl")
        from openpyxl import Workbook

        wb = Workbook()
        ws = wb.active
        ws.title = "Sheet1"
        ws.append(["Name", "Value"])
        ws.append(["Alpha", 100])
        f = tmp_path / "test.xlsx"
        wb.save(str(f))
        wb.close()

        result = extract_text_xlsx(f)
        assert result is not None
        assert "[Sheet1]" in result
        assert "Name | Value" in result
        assert "Alpha | 100" in result


class TestExtractTextDOCX:
    def test_extract_docx(self, tmp_path: Path) -> None:
        """DOCX files return paragraph text."""
        pytest.importorskip("docx")
        from docx import Document

        doc = Document()
        doc.add_heading("Test Document", level=1)
        doc.add_paragraph("First paragraph.")
        doc.add_paragraph("Second paragraph.")
        f = tmp_path / "test.docx"
        doc.save(str(f))

        result = extract_text_docx(f)
        assert result is not None
        assert "Test Document" in result
        assert "First paragraph." in result
        assert "Second paragraph." in result


class TestExtractTextPPTX:
    def test_extract_pptx(self, tmp_path: Path) -> None:
        """PPTX files return slide text with headers."""
        pytest.importorskip("pptx")
        from pptx import Presentation

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Slide Title"
        f = tmp_path / "test.pptx"
        prs.save(str(f))

        result = extract_text_pptx(f)
        assert result is not None
        assert "[Slide 1]" in result
        assert "Slide Title" in result


class TestExtractTextFromRouter:
    def test_routes_to_plain(self, tmp_path: Path) -> None:
        """Router dispatches .txt to plain extractor."""
        f = tmp_path / "doc.txt"
        f.write_text("content here", encoding="utf-8")
        result = extract_text_from_file(f)
        assert result == "content here"

    def test_routes_to_csv(self, tmp_path: Path) -> None:
        """Router dispatches .csv to CSV extractor."""
        f = tmp_path / "data.csv"
        f.write_text("a,b\n1,2", encoding="utf-8")
        result = extract_text_from_file(f)
        assert result is not None
        assert "a | b" in result

    def test_unsupported_returns_none(self, tmp_path: Path) -> None:
        """Unsupported extension returns None."""
        f = tmp_path / "image.xyz"
        f.write_text("data", encoding="utf-8")
        result = extract_text_from_file(f)
        assert result is None

    def test_document_extensions_list(self) -> None:
        """DOCUMENT_EXTENSIONS includes all expected types."""
        expected = {'.pdf', '.docx', '.doc', '.txt', '.md', '.rtf',
                    '.xlsx', '.csv', '.pptx'}
        assert expected.issubset(set(DOCUMENT_EXTENSIONS))
