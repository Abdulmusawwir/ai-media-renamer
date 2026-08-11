from __future__ import annotations

import base64
import datetime
import functools
import hashlib
import json
import logging
import os
import shutil
import subprocess
import sys
from abc import ABC, abstractmethod
from pathlib import Path
from typing import Any, Callable

import anthropic
import keyring
import ollama
import openai
import requests
from pydantic import BaseModel, ValidationError

VERSION = "v1.6.2"


class AssetAnalysisResponse(BaseModel):
    """Structured Pydantic model for AI analysis responses.

    Defines the expected JSON schema returned by all AI providers.
    Used for validation, retry logic, and optional provider-level
    structured output (response_format).
    """

    new_filename: str
    suggested_category: str
    overall_visual_summary: str
    tags: list[str]
    topic: str = ""
    description: str = ""
    confidence: float = 0.0


_NO_WINDOW = getattr(subprocess, "CREATE_NO_WINDOW", 0)

# -----------------------------------------------------------------------------
# 0. HELPER: resolve worker count
# -----------------------------------------------------------------------------

def _resolve_workers(cfg_value: int | float | None) -> int:
    """Resolve worker count from config value, falling back to CPU count.

    Args:
        cfg_value: Configured worker count from config, or None/0 for auto.

    Returns:
        Positive integer worker count.
    """
    if isinstance(cfg_value, int) and cfg_value > 0:
        return cfg_value
    return os.cpu_count() or 4

# -----------------------------------------------------------------------------
# 1. CONFIGURATION & LOGGING
# -----------------------------------------------------------------------------

def load_config(config_path: str = "config.json") -> dict[str, Any]:
    """Load configuration from JSON file with auto-recovery from default.

    Args:
        config_path: Relative path to the config file from the script directory.

    Returns:
        Parsed configuration dictionary with normalized tuple fields.
    """
    script_dir = Path(__file__).parent
    full_path = script_dir / config_path
    default_path = full_path.parent / "config.default.json"
    try:
        with open(full_path, encoding='utf-8') as f:
            cfg = json.load(f)
    except (FileNotFoundError, json.JSONDecodeError) as e:
        print(f"Warning: config.json issue: {e}. Attempting auto-recovery from config.default.json...")
        try:
            import shutil
            shutil.copy2(default_path, full_path)
            with open(full_path, encoding='utf-8') as f:
                cfg = json.load(f)
            print(f"Auto-recovery successful: restored config.json from config.default.json")
        except Exception as recovery_err:
            print(f"Error: Auto-recovery failed: {recovery_err}")
            print(f"To fix: manually copy config.default.json to config.json, or run with --reset-config")
            sys.exit(1)

    cfg['video_extensions'] = tuple(cfg.get('video_extensions', ['.mp4', '.mov', '.avi', '.mkv', '.webm']))
    cfg['image_extensions'] = tuple(cfg.get('image_extensions', ['.jpg', '.jpeg', '.png', '.webp', '.gif']))
    cfg['audio_extensions'] = tuple(cfg.get('audio_extensions', [
        '.mp3', '.wav', '.flac', '.aac', '.ogg', '.m4a',
        '.wma', '.opus', '.aiff', '.alac', '.ape', '.wv',
    ]))
    cfg['allowed_categories'] = tuple(cfg.get('allowed_categories', []))
    return cfg


def restore_default_config() -> bool:
    """Overwrite config.json with config.default.json. Returns True on success."""
    import shutil
    script_dir = Path(__file__).parent
    default_path = script_dir / "config.default.json"
    target_path = script_dir / "config.json"
    if not default_path.exists():
        return False
    shutil.copy2(default_path, target_path)
    return True


config = load_config()

ALLOWED_CATEGORIES = config['allowed_categories']


def get_active_profile() -> str:
    """Return the name of the currently active prompt profile."""
    return config.get('prompt_profiles', {}).get('active', 'general_balanced')


def get_active_categories() -> tuple[str, ...]:
    """Return the allowed categories for the active prompt profile."""
    profile_name = get_active_profile()
    profile = config.get('prompt_profiles', {}).get('profiles', {}).get(profile_name, {})
    cats = profile.get('allowed_categories', [])
    return tuple(cats) if cats else ALLOWED_CATEGORIES


def get_active_prompt() -> str:
    """Return the active prompt text with categories expanded into the template."""
    profile_name = get_active_profile()
    profile = config.get('prompt_profiles', {}).get('profiles', {}).get(profile_name, {})
    raw = profile.get('prompt', '')
    cats = get_active_categories()
    cat_str = "\n".join(f'   - "{c}"' for c in cats)

    for needle in ("the allowed categories list", "the allowed list"):
        if needle in raw:
            raw = raw.replace(needle, f"this list:\n{cat_str}")
            break

    constraint = (
        f"\n\nIMPORTANT — ALLOWED CATEGORIES (use ONLY these, never invent new ones):\n{cat_str}"
    )
    return raw + constraint


def set_active_profile(name: str) -> None:
    """Set the active prompt profile by name and persist to config.

    Args:
        name: Profile key to activate.
    """
    profiles = config.get('prompt_profiles', {}).get('profiles', {})
    if name in profiles:
        config['prompt_profiles']['active'] = name
        save_config()


def get_profile_labels() -> dict[str, str]:
    """Return a mapping of profile keys to their display labels."""
    profiles = config.get('prompt_profiles', {}).get('profiles', {})
    return {k: v.get('label', k) for k, v in profiles.items()}


PROMPT_PROFILES = get_profile_labels()

VIDEO_EXTENSIONS = config['video_extensions']
IMAGE_EXTENSIONS = config['image_extensions']
AUDIO_EXTENSIONS = config.get('audio_extensions', (
    '.mp3', '.wav', '.flac', '.aac', '.ogg', '.m4a',
    '.wma', '.opus', '.aiff', '.alac', '.ape', '.wv',
))
MODEL_NAME = config['model']['name']
TEXT_MODEL_NAME = config['model'].get('text_model', 'qwen2.5:3b')
MODEL_TEMPERATURE = config['model']['temperature']
MODEL_NUM_CTX = config['model']['num_ctx']
MODEL_KEEP_ALIVE = config['model']['keep_alive']
IMAGE_PREVIEW_MAX_EDGE = config['preview']['image_max_edge']
VIDEO_GRID_SCALE = config['preview']['video_grid_scale']
EXTRACTION_WORKERS = _resolve_workers(config['preview'].get('extraction_workers', 0))

DEFAULT_CASE_STYLE = config.get('naming', {}).get('case_style', 'title_case')
DEFAULT_MAX_FILENAME_CHARS = config.get('naming', {}).get('max_filename_chars', 0)

CURRENT_PROVIDER = config.get('model', {}).get('last_provider', 'ollama')
CURRENT_API_KEY = ""

NAMED_TEMPLATES = config.get('naming_templates', {
    "default": "{topic}_{description}",
    "short": "{topic}_{description}",
    "editorial": "{date}_{topic}"
})
DEFAULT_TEMPLATE_STRING = NAMED_TEMPLATES.get("default", "{topic}_{description}")

LOG_DIR = Path(config['logging']['directory'])
MAX_UPLOAD_SIZE = int(config['logging'].get('max_upload_size', 10737418240))
CONFIG_PATH = Path(__file__).parent / "config.json"
KEYRING_SERVICE = "ai-media-renamer"
PROVIDER_REGISTRY = {}


def save_config() -> None:
    """Persist the current in-memory config dict to config.json."""
    global config
    with open(CONFIG_PATH, 'w', encoding='utf-8') as f:
        json.dump(config, f, indent=2)


def reload_config() -> None:
    """Reload config from disk and refresh all module-level globals."""
    global config, ALLOWED_CATEGORIES, VIDEO_EXTENSIONS, IMAGE_EXTENSIONS, AUDIO_EXTENSIONS
    global MODEL_NAME, MODEL_TEMPERATURE, MODEL_NUM_CTX, MODEL_KEEP_ALIVE
    global EXTRACTION_WORKERS, DEFAULT_CASE_STYLE, DEFAULT_MAX_FILENAME_CHARS
    global NAMED_TEMPLATES, DEFAULT_TEMPLATE_STRING, PROMPT_PROFILES, CURRENT_PROVIDER
    config = load_config()
    ALLOWED_CATEGORIES = config['allowed_categories']
    VIDEO_EXTENSIONS = config['video_extensions']
    IMAGE_EXTENSIONS = config['image_extensions']
    AUDIO_EXTENSIONS = config['audio_extensions']
    MODEL_NAME = config['model']['name']
    MODEL_TEMPERATURE = config['model']['temperature']
    MODEL_NUM_CTX = config['model']['num_ctx']
    MODEL_KEEP_ALIVE = config['model']['keep_alive']
    EXTRACTION_WORKERS = _resolve_workers(config['preview'].get('extraction_workers', 0))
    DEFAULT_CASE_STYLE = config.get('naming', {}).get('case_style', 'title_case')
    DEFAULT_MAX_FILENAME_CHARS = config.get('naming', {}).get('max_filename_chars', 0)
    NAMED_TEMPLATES = config.get('naming_templates', {
        "default": "{topic}_{description}",
        "short": "{topic}_{description}",
        "editorial": "{date}_{topic}"
    })
DEFAULT_TEMPLATE_STRING = NAMED_TEMPLATES.get("default", "{topic}_{description}")

# -----------------------------------------------------------------------------
# 1b. SETUP / ONBOARDING (use-case → one-time dependency matrix)
# -----------------------------------------------------------------------------

USER_DATA_DIR = Path(os.environ.get("APPDATA", Path.home() / "AppData" / "Roaming")) / "ai-media-renamer"

# Use-case keys map to the external tools + model kinds they require. Kept in
# engine.py so both bootstrap.py (setup wizard) and app.py (env checks) share
# a single source of truth.
SETUP_USE_CASES: dict[str, dict[str, Any]] = {
    "videos": {
        "label": "Videos",
        "desc": "MP4, MOV, MKV, AVI \u2014 scene analysis with audio-track context",
        "needs": {"ffmpeg", "exiftool", "vision_model", "whisper"},
    },
    "photos": {
        "label": "Photos",
        "desc": "JPG, PNG, WEBP, GIF \u2014 visual analysis and tagging",
        "needs": {"ffmpeg", "exiftool", "vision_model"},
    },
    "documents": {
        "label": "Documents",
        "desc": "PDF, DOCX, TXT, MD, RTF \u2014 content-based renaming",
        "needs": {"text_model"},
    },
    "spreadsheets": {
        "label": "Spreadsheets",
        "desc": "XLSX, CSV \u2014 data-aware renaming",
        "needs": {"text_model"},
    },
    "audio": {
        "label": "Audio",
        "desc": "MP3, WAV, FLAC \u2014 transcription + content renaming",
        "needs": {"exiftool", "text_model", "whisper"},
    },
}

SETUP_DEPENDENCIES: dict[str, dict[str, Any]] = {
    "ollama": {
        "label": "Ollama runtime",
        "desc": "Runs your AI models locally \u2014 required for every use case.",
        "always": True,
    },
    "ffmpeg": {
        "label": "FFmpeg",
        "desc": "Extracts video frames and downscales images.",
    },
    "exiftool": {
        "label": "ExifTool",
        "desc": "Writes metadata tags (needed for PDF metadata too).",
    },
    "vision_model": {
        "label": "Vision AI model",
        "desc": "Analyzes video frames and photos.",
    },
    "text_model": {
        "label": "Text AI model",
        "desc": "Analyzes documents, spreadsheets, and transcripts.",
    },
    "whisper": {
        "label": "Whisper (speech-to-text)",
        "desc": "Transcribes audio tracks and audio files.",
    },
}

# Model catalog shared by the setup wizard and the in-app model dropdowns.
# `kind` selects the category; recommended_* flags pick a default per hardware.
MODEL_CATALOG: list[dict[str, Any]] = [
    # ---- vision models ----
    {
        "name": "qwen2.5vl:7b", "label": "Qwen 2.5 VL 7B", "kind": "vision",
        "size": "6.0 GB", "size_gb": 6.0, "quality": "Best", "speed": "Moderate",
        "desc": "Best quality for structured JSON extraction and visual analysis.",
        "recommended_gpu": True,
    },
    {
        "name": "qwen2.5vl:3b", "label": "Qwen 2.5 VL 3B", "kind": "vision",
        "size": "3.2 GB", "size_gb": 3.2, "quality": "Good", "speed": "Fast",
        "desc": "Lighter model \u2014 good quality, faster on CPU or low VRAM.",
        "recommended_cpu": True,
    },
    {
        "name": "qwen3-vl:4b", "label": "Qwen 3 VL 4B", "kind": "vision",
        "size": "~3 GB", "size_gb": 3.0, "quality": "Good", "speed": "Fast",
        "desc": "Newer architecture with improved reasoning.",
    },
    {
        "name": "moondream:latest", "label": "Moondream 2", "kind": "vision",
        "size": "1.8 GB", "size_gb": 1.8, "quality": "Basic", "speed": "Very fast",
        "desc": "Smallest option \u2014 basic visual understanding for very low VRAM.",
    },
    # ---- text models ----
    {
        "name": "qwen2.5:3b", "label": "Qwen 2.5 3B", "kind": "text",
        "size": "1.9 GB", "size_gb": 1.9, "quality": "Good", "speed": "Fast",
        "desc": "Recommended text model \u2014 fast on CPU, great for documents and transcripts.",
        "recommended_cpu": True, "recommended_gpu": True,
    },
    {
        "name": "qwen2.5:7b", "label": "Qwen 2.5 7B", "kind": "text",
        "size": "4.7 GB", "size_gb": 4.7, "quality": "Best", "speed": "Moderate",
        "desc": "Highest text quality \u2014 slower on CPU.",
        "recommended_gpu": True,
    },
    {
        "name": "qwen2.5:1.5b", "label": "Qwen 2.5 1.5B", "kind": "text",
        "size": "986 MB", "size_gb": 0.99, "quality": "Basic", "speed": "Very fast",
        "desc": "Lightest option for very old hardware.",
    },
]

# GGUF models downloaded by the setup wizard for the llama.cpp runtime. These
# are raw model files served by llama-server; a vision entry carries an mmproj
# (vision-projector) companion file and ALSO handles text-only prompts, so a
# profile needing both vision and text only downloads the vision GGUF.
LLAMACPP_GGUF_CATALOG: list[dict[str, Any]] = [
    {
        "name": "qwen2.5vl:7b", "kind": "vision", "label": "Qwen 2 VL 7B (Q4_K_M)",
        "size": "5.4 GB", "size_gb": 5.4, "quality": "Best", "speed": "Moderate",
        "desc": "Vision + text in one file. Best quality for structured JSON "
                "extraction and visual analysis.",
        "url": "https://huggingface.co/ggml-org/Qwen2-VL-7B-Instruct-GGUF/resolve/main/"
               "Qwen2-VL-7B-Instruct-Q4_K_M.gguf",
        "sha256": "cba46c253ee6d1bd4c322f5a620cf69d65d644086c12df8dba59ebeff0768501",
        "mmproj_url": "https://huggingface.co/ggml-org/Qwen2-VL-7B-Instruct-GGUF/resolve/main/"
                      "mmproj-Qwen2-VL-7B-Instruct-Q8_0.gguf",
        "mmproj_sha256": "97fbf5ee6c08b6fb34b9d589d2531d980714401b6150db6ee716fcf45b215bc4",
        "recommended_gpu": True,
    },
    {
        "name": "qwen2.5vl:2b", "kind": "vision", "label": "Qwen 2 VL 2B (Q4_K_M)",
        "size": "1.7 GB", "size_gb": 1.7, "quality": "Good", "speed": "Fast",
        "desc": "Lighter option \u2014 good quality, fast on CPU or low VRAM.",
        "url": "https://huggingface.co/ggml-org/Qwen2-VL-2B-Instruct-GGUF/resolve/main/"
               "Qwen2-VL-2B-Instruct-Q4_K_M.gguf",
        "sha256": "5745685d2e607a82a0696c1118e56a2a1ae0901da450fd9cd4f161c6b62867d7",
        "mmproj_url": "https://huggingface.co/ggml-org/Qwen2-VL-2B-Instruct-GGUF/resolve/main/"
                      "mmproj-Qwen2-VL-2B-Instruct-Q8_0.gguf",
        "mmproj_sha256": "a0ad91f00a7a80dcf84d719a61b00ee2e07b71794f4ee2dfa81a254621a8c418",
        "recommended_cpu": True,
    },
    {
        "name": "qwen2.5:3b", "kind": "text", "label": "Qwen 2.5 3B (Q4_K_M)",
        "size": "2.1 GB", "size_gb": 2.1, "quality": "Good", "speed": "Fast",
        "desc": "Text-only model \u2014 recommended for documents, spreadsheets "
                "and transcripts.",
        "url": "https://huggingface.co/Qwen/Qwen2.5-3B-Instruct-GGUF/resolve/main/"
               "qwen2.5-3b-instruct-q4_k_m.gguf",
        "sha256": "626b4a6678b86442240e33df819e00132d3ba7dddfe1cdc4fbb18e0a9615c62d",
        "mmproj_url": "",
        "mmproj_sha256": "",
        "recommended_cpu": True, "recommended_gpu": True,
    },
]


def use_cases_needs(profile: list[str]) -> set[str]:
    """Return the set of dependencies required for the given use cases."""
    needs: set[str] = set()
    for key in profile:
        needs |= SETUP_USE_CASES.get(key, {}).get("needs", set())
    return needs


def _has_gpu() -> bool:
    """Cheap GPU detection for setup recommendations (works before FFmpeg exists)."""
    try:
        result = subprocess.run(
            ["nvidia-smi", "--query-gpu=name", "--format=csv,noheader"],
            capture_output=True, text=True, timeout=5, creationflags=_NO_WINDOW,
        )
        if result.returncode == 0 and result.stdout.strip():
            return True
    except Exception:
        pass
    if sys.platform == "win32":
        import ctypes
        for dll in ("amdx64.dll", "atiadlxx.dll"):
            try:
                ctypes.windll.LoadLibrary(dll)
                return True
            except Exception:
                continue
    return False


def recommended_models(profile: list[str]) -> dict[str, str]:
    """Pick default vision/text models for a use-case profile + detected hardware.

    Returns:
        Dict with optional 'vision' and 'text' model names.
    """
    needs = use_cases_needs(profile)
    has_gpu = _has_gpu()
    out: dict[str, str] = {}
    for kind in ("vision", "text"):
        if f"{kind}_model" not in needs:
            continue
        cands = [m for m in MODEL_CATALOG if m["kind"] == kind]
        preferred = next(
            (m for m in cands if m.get("recommended_gpu" if has_gpu else "recommended_cpu")),
            None,
        )
        out[kind] = (preferred or cands[0])["name"]
    return out


def recommended_llamacpp_models(profile: list[str]) -> dict[str, str]:
    """Pick default GGUF vision/text models for the llama.cpp runtime.

    Returns:
        Dict with optional 'vision' and 'text' model names (from
        LLAMACPP_GGUF_CATALOG), chosen for the detected hardware.
    """
    needs = use_cases_needs(profile)
    has_gpu = _has_gpu()
    out: dict[str, str] = {}
    for kind in ("vision", "text"):
        if f"{kind}_model" not in needs:
            continue
        cands = [m for m in LLAMACPP_GGUF_CATALOG if m["kind"] == kind]
        if not cands:
            continue
        preferred = next(
            (m for m in cands if m.get("recommended_gpu" if has_gpu else "recommended_cpu")),
            None,
        )
        out[kind] = (preferred or cands[0])["name"]
    return out


@functools.lru_cache(maxsize=128)
def _model_tag_exists(name: str, tag: str) -> bool | None:
    """True/False if the tag exists in Ollama's registry, None if offline."""
    try:
        resp = requests.get(
            f"https://registry.ollama.ai/v2/library/{name}/tags/list", timeout=10
        )
        if resp.status_code != 200:
            return None
        data = resp.json()
        tags = [t.get("name", t) if isinstance(t, dict) else t for t in data.get("tags", [])]
        return tag in tags
    except Exception:
        return None


def validate_ollama_model(model_name: str) -> bool | None:
    """Confirm a model tag exists on Ollama's registry (None when offline)."""
    name, _, tag = model_name.partition(":")
    return _model_tag_exists(name, tag or "latest")


def pre_download_whisper(model_size: str = "base") -> bool:
    """Pre-download a faster-whisper model so the first transcription is instant."""
    try:
        from huggingface_hub import snapshot_download
        snapshot_download(repo_id=f"Systran/faster-whisper-{model_size}")
        return True
    except Exception:
        return False


def load_setup_profile() -> dict[str, Any]:
    """Load the persisted onboarding profile from user data (%APPDATA%)."""
    default: dict[str, Any] = {"onboarded": False, "profile": [], "setup_version": VERSION,
                               "runtime": ""}
    try:
        data = json.loads((USER_DATA_DIR / "setup.json").read_text(encoding="utf-8"))
        data.setdefault("onboarded", False)
        data.setdefault("profile", [])
        data.setdefault("setup_version", VERSION)
        data.setdefault("runtime", "")
        return data
    except Exception:
        return default


def save_setup_profile(profile: list[str] | None = None, onboarded: bool = True,
                       runtime: str | None = None) -> dict[str, Any]:
    """Persist the onboarding profile to user data so it survives EXE restarts.

    Args:
        profile: Use-case keys chosen in onboarding ('', to clear).
        onboarded: True once the questionnaire has been answered.
        runtime: Chosen local AI runtime ('ollama' or 'llamacpp').
    """
    data = load_setup_profile()
    if profile is not None:
        data["profile"] = list(profile)
    data["onboarded"] = onboarded
    data["setup_version"] = VERSION
    if runtime is not None:
        data["runtime"] = runtime
    try:
        USER_DATA_DIR.mkdir(parents=True, exist_ok=True)
        (USER_DATA_DIR / "setup.json").write_text(json.dumps(data, indent=2), encoding="utf-8")
    except Exception:
        pass
    return data
    PROMPT_PROFILES = get_profile_labels()
    CURRENT_PROVIDER = config.get('model', {}).get('last_provider', 'ollama')


def save_api_key(provider_name: str, key: str) -> None:
    """Store an API key for a provider in the system keyring.

    Args:
        provider_name: Provider identifier (e.g. 'gemini', 'openai').
        key: API key string to store.
    """
    keyring.set_password(KEYRING_SERVICE, provider_name, key)


def load_api_key(provider_name: str) -> str:
    """Retrieve an API key from the system keyring.

    Args:
        provider_name: Provider identifier to look up.

    Returns:
        The stored API key, or an empty string if not found.
    """
    return keyring.get_password(KEYRING_SERVICE, provider_name) or ""


def delete_api_key(provider_name: str) -> None:
    """Delete an API key from the system keyring, ignoring if absent.

    Args:
        provider_name: Provider identifier whose key to delete.
    """
    try:
        keyring.delete_password(KEYRING_SERVICE, provider_name)
    except keyring.errors.PasswordDeleteError:
        pass


def setup_logging(verbose: bool = False) -> logging.Logger:
    """Configure and return the application JSONL file logger.

    Args:
        verbose: If True, set log level to DEBUG; otherwise INFO.

    Returns:
        Configured logger instance writing to the daily log file.
    """
    log_dir = LOG_DIR
    log_dir.mkdir(parents=True, exist_ok=True)
    log_file = log_dir / f"renamer_{datetime.datetime.now().astimezone().date().isoformat()}.jsonl"

    logger = logging.getLogger('video_renamer')
    logger.handlers.clear()
    logger.setLevel(logging.DEBUG if verbose else logging.INFO)

    file_handler = logging.FileHandler(log_file, encoding='utf-8', mode='a')
    file_handler.setLevel(logging.DEBUG if verbose else logging.INFO)
    logger.addHandler(file_handler)

    return logger


def log_event(logger: Any, level: str, event: str, file_name: str | None = None, details: dict[str, Any] | None = None) -> None:
    """Write a structured JSON log entry.

    Args:
        logger: Logger instance to write to.
        level: Log level string ('DEBUG', 'INFO', 'WARNING', 'ERROR').
        event: Short event description.
        file_name: Optional name of the file being processed.
        details: Optional dictionary of additional context.
    """
    record = {
        "timestamp": datetime.datetime.now().astimezone().isoformat(),
        "level": level,
        "event": event,
        "file": file_name,
    }
    if details:
        record["details"] = details
    msg = json.dumps(record)
    if level == "DEBUG":
        logger.debug(msg)
    elif level == "WARNING":
        logger.warning(msg)
    elif level == "ERROR":
        logger.error(msg)
    else:
        logger.info(msg)


# -----------------------------------------------------------------------------
# 2. EXIFTOOL PERSISTENT BACKGROUND PROCESS (stay_open)
# -----------------------------------------------------------------------------

class ExifToolSession:
    def __init__(self) -> None:
        """Start a persistent ExifTool subprocess in stay_open mode."""
        try:
            self.process = subprocess.Popen(
                ['exiftool', '-stay_open', 'True', '-@', '-'],
                stdin=subprocess.PIPE, stdout=subprocess.PIPE, stderr=subprocess.STDOUT,
                text=True, encoding='utf-8', bufsize=1,
                creationflags=_NO_WINDOW,
            )
        except FileNotFoundError:
            print("Error: ExifTool is not installed or not in system PATH.")
            sys.exit(1)

    def execute(self, args: list[str]) -> str:
        """Send arguments to ExifTool and return the output.

        Args:
            args: List of ExifTool arguments (flags and file paths).

        Returns:
            Raw text output from ExifTool.
        """
        for arg in args:
            self.process.stdin.write(f"{arg}\n")
        self.process.stdin.write("-execute\n")
        self.process.stdin.flush()

        output = ""
        for line in self.process.stdout:
            if "{ready}" in line:
                break
            output += line
        return output

    def execute_batch(self, all_args: list[list[str]]) -> list[str]:
        """Send multiple file argument sets in a single -execute block.

        Batches all metadata writes into one IPC round-trip instead of one
        per file, reducing overhead from ~200ms x N to ~200ms total.

        Args:
            all_args: List of argument lists, one per file. Each inner list
                      ends with the file path as the last element.

        Returns:
            List of ExifTool output strings, one per file argument set.
        """
        for file_args in all_args:
            for arg in file_args:
                self.process.stdin.write(f"{arg}\n")
        self.process.stdin.write("-execute\n")
        self.process.stdin.flush()

        output = ""
        for line in self.process.stdout:
            if "{ready}" in line:
                break
            output += line
        return [output]

    def close(self) -> None:
        """Shut down the persistent ExifTool subprocess."""
        if hasattr(self, 'process'):
            self.process.stdin.write("-stay_open\nFalse\n")
            self.process.stdin.flush()
            self.process.wait()


# -----------------------------------------------------------------------------
# 3. HARDWARE & CACHE MANAGERS
# -----------------------------------------------------------------------------

def detect_hw_accel() -> str | None:
    """Probe for available hardware-accelerated FFmpeg decoders.

    Returns:
        Hardware accelerator name ('cuda', 'qsv', 'amf') or None.
    """
    for hw in ['cuda', 'qsv', 'amf']:
        try:
            cmd = ['ffmpeg', '-hwaccel', hw, '-f', 'lavfi', '-i', 'color=c=black:s=16x16:d=1', '-f', 'null', '-']
            res = subprocess.run(cmd, capture_output=True, creationflags=_NO_WINDOW)
            if res.returncode == 0:
                return hw
        except Exception:
            pass
    return None


def is_already_processed(file_path: str | Path, exiftool_session: Any) -> bool:
    """Check if a file already has XMP Description metadata written.

    Args:
        file_path: Path to the media file.
        exiftool_session: Active ExifToolSession instance.

    Returns:
        True if the file already has a DC:Description tag.
    """
    output = exiftool_session.execute(["-XMP-dc:Description", "-json", str(file_path)])
    try:
        data = json.loads(output.strip())
        if isinstance(data, list) and len(data) > 0:
            return "XMP-dc:Description" in data[0]
    except (json.JSONDecodeError, TypeError, IndexError):
        pass
    return False


# -----------------------------------------------------------------------------
# 4. ZERO-I/O PIPELINE (MEMORY-BASED ASSET EXTRACTION)
# -----------------------------------------------------------------------------

def get_video_duration(video_path: str | Path) -> float:
    """Return the duration in seconds of a video file, with internal caching.

    Args:
        video_path: Path to the video file.

    Returns:
        Duration in seconds, defaulting to 10.0 on probe failure.
    """
    cache_key = str(video_path)
    if not hasattr(get_video_duration, '_cache'):
        get_video_duration._cache = {}
    if cache_key in get_video_duration._cache:
        return get_video_duration._cache[cache_key]
    cmd = ['ffprobe', '-v', 'error', '-show_entries', 'format=duration',
           '-of', 'default=noprint_wrappers=1:nokey=1', str(video_path)]
    try:
        output = subprocess.check_output(cmd, stderr=subprocess.STDOUT, creationflags=_NO_WINDOW).decode().strip()
        duration = float(output)
    except Exception:
        duration = 10.0
    get_video_duration._cache[cache_key] = duration
    return duration


def _extract_frame_to_base64(video_path: str | Path, hw_accel: str | None) -> str | None:
    """Extract the midpoint frame of a video and return as base64 JPEG.

    Args:
        video_path: Path to the video file.
        hw_accel: Hardware accelerator name or None for software decoding.

    Returns:
        Base64-encoded JPEG string, or None on failure.
    """
    duration = get_video_duration(video_path)
    mid_offset = max(1.0, duration * 0.5)

    cmd = ['ffmpeg', '-y', '-hide_banner', '-loglevel', 'error']
    if hw_accel:
        cmd.extend(['-hwaccel', hw_accel])

    cmd.extend([
        '-ss', str(mid_offset),
        '-i', str(video_path),
        '-vframes', '1',
        '-vf', f"scale={VIDEO_GRID_SCALE}:-1",
        '-f', 'image2pipe',
        '-vcodec', 'mjpeg',
        '-'
    ])

    try:
        process = subprocess.run(cmd, capture_output=True, check=True, creationflags=_NO_WINDOW)
        return base64.b64encode(process.stdout).decode('utf-8')
    except (subprocess.CalledProcessError, FileNotFoundError):
        return None


def process_video_to_base64(video_path: str | Path, hw_accel: str | None,
                            fallback_log: dict[str, bool] | None = None) -> str | None:
    """Extract the midpoint frame of a video, falling back to CPU on failure.

    If hardware-accelerated decoding is requested but fails, retries the
    extraction with software decoding and records the file name in
    ``fallback_log`` so callers can notify the user about the degradation.

    Args:
        video_path: Path to the video file.
        hw_accel: Hardware accelerator name or None for software decoding.
        fallback_log: Optional dict; file names that needed CPU fallback are
            recorded here (name -> True).

    Returns:
        Base64-encoded JPEG string, or None on failure.
    """
    if hw_accel:
        result = _extract_frame_to_base64(video_path, hw_accel)
        if result is not None:
            return result
        if fallback_log is not None:
            fallback_log[Path(video_path).name] = True
    return _extract_frame_to_base64(video_path, None)


def process_image_to_base64(image_path: str | Path, max_edge: int = IMAGE_PREVIEW_MAX_EDGE) -> str | None:
    """Downscale an image and return as base64 JPEG.

    Args:
        image_path: Path to the image file.
        max_edge: Maximum pixel dimension for the longest edge.

    Returns:
        Base64-encoded JPEG string, or None on failure.
    """
    cmd = [
        'ffmpeg', '-y', '-hide_banner', '-loglevel', 'error',
        '-i', str(image_path),
        '-vf', f"scale={max_edge}:{max_edge}:force_original_aspect_ratio=decrease",
        '-frames:v', '1',
        '-f', 'image2pipe',
        '-vcodec', 'mjpeg',
        '-q:v', '3',
        '-'
    ]
    try:
        process = subprocess.run(cmd, capture_output=True, check=True, creationflags=_NO_WINDOW)
        return base64.b64encode(process.stdout).decode('utf-8')
    except subprocess.CalledProcessError:
        return None
    except FileNotFoundError:
        return None


def process_asset_to_base64(file_path: Path, hw_accel: str | None,
                            fallback_log: dict[str, bool] | None = None) -> str | None:
    """Route a media file to the appropriate base64 encoder.

    Args:
        file_path: Path to the video or image file.
        hw_accel: Hardware accelerator name or None.
        fallback_log: Optional dict forwarded to ``process_video_to_base64``
            to record files that fell back to CPU decoding.

    Returns:
        Base64-encoded JPEG string, or None on failure.
    """
    if file_path.suffix.lower() in VIDEO_EXTENSIONS:
        return process_video_to_base64(file_path, hw_accel, fallback_log)
    return process_image_to_base64(file_path)


# -----------------------------------------------------------------------------
# 3c. AUDIO TRANSCRIPTION
# -----------------------------------------------------------------------------


def extract_audio_from_video(video_path: str | Path) -> Path | None:
    """Extract audio track from a video file using FFmpeg.

    Outputs 16kHz mono WAV suitable for Whisper transcription.

    Args:
        video_path: Path to the video file.

    Returns:
        Path to extracted WAV file, or None if no audio track or on error.
    """
    import tempfile
    video_path = Path(video_path)
    if not video_path.exists():
        return None
    tmp = Path(tempfile.mktemp(suffix=".wav"))
    cmd = [
        'ffmpeg', '-y', '-hide_banner', '-loglevel', 'error',
        '-i', str(video_path),
        '-vn', '-acodec', 'pcm_s16le', '-ar', '16000', '-ac', '1',
        str(tmp),
    ]
    try:
        subprocess.run(cmd, capture_output=True, check=True, creationflags=_NO_WINDOW)
        return tmp if tmp.exists() else None
    except (subprocess.CalledProcessError, FileNotFoundError):
        if tmp.exists():
            tmp.unlink()
        return None


def _has_audio_track(video_path: str | Path) -> bool:
    """Check if a video file has an audio track via FFprobe.

    Args:
        video_path: Path to the video file.

    Returns:
        True if the video has at least one audio stream.
    """
    cmd = [
        'ffprobe', '-v', 'error',
        '-select_streams', 'a',
        '-show_entries', 'stream=index',
        '-of', 'csv=p=0',
        str(video_path),
    ]
    try:
        result = subprocess.run(cmd, capture_output=True, text=True, creationflags=_NO_WINDOW)
        return bool(result.stdout.strip())
    except (subprocess.CalledProcessError, FileNotFoundError):
        return False


_whisper_model_cache: dict[str, Any] = {}


def transcribe_audio(audio_path: str | Path, model_size: str = "base") -> dict[str, Any]:
    """Transcribe an audio file using faster-whisper (local, no cloud).

    The model is lazily loaded and cached in memory for subsequent calls.

    Args:
        audio_path: Path to the audio file (WAV, MP3, etc.).
        model_size: Whisper model size — tiny (39MB), base (74MB),
                    small (244MB), medium (769MB), large-v3 (1.5GB).

    Returns:
        Dict with keys: text (str), language (str), duration (float).
        On error: text is empty and 'error' key is set.
    """
    audio_path = Path(audio_path)
    if not audio_path.exists():
        return {"text": "", "language": "", "duration": 0.0, "error": "File not found"}

    if model_size not in _whisper_model_cache:
        try:
            from faster_whisper import WhisperModel
            _whisper_model_cache[model_size] = WhisperModel(model_size, device="cpu", compute_type="int8")
        except ImportError:
            return {"text": "", "language": "", "duration": 0.0,
                    "error": "faster-whisper not installed. Run: pip install faster-whisper"}
        except Exception as exc:
            return {"text": "", "language": "", "duration": 0.0,
                    "error": f"Failed to load whisper model: {exc}"}

    model = _whisper_model_cache[model_size]
    try:
        segments, info = model.transcribe(str(audio_path), beam_size=1)
        text_parts = [seg.text for seg in segments]
        return {
            "text": " ".join(text_parts).strip(),
            "language": info.language or "",
            "duration": info.duration or 0.0,
        }
    except Exception as exc:
        return {"text": "", "language": "", "duration": 0.0, "error": str(exc)}


# -----------------------------------------------------------------------------
# 4b. DOCUMENT TEXT EXTRACTION
# -----------------------------------------------------------------------------

DOCUMENT_EXTENSIONS = config.get('document_extensions', [
    '.pdf', '.docx', '.doc', '.txt', '.md', '.rtf',
    '.xlsx', '.csv', '.pptx',
])

SPREADSHEET_EXTENSIONS = {'.xlsx', '.csv'}
PRESENTATION_EXTENSIONS = {'.pptx'}


def extract_text_pdf(path: Path) -> str | None:
    """Extract text from a PDF file using pdfplumber.

    Args:
        path: Path to the PDF file.

    Returns:
        Extracted text, or None on failure.
    """
    try:
        import logging
        logging.getLogger("pdfminer").setLevel(logging.ERROR)
        import pdfplumber
        text_parts: list[str] = []
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(page_text)
        return "\n\n".join(text_parts).strip() or None
    except Exception:
        return None


def extract_text_docx(path: Path) -> str | None:
    """Extract text from a DOCX file using python-docx.

    Args:
        path: Path to the DOCX file.

    Returns:
        Extracted text, or None on failure.
    """
    try:
        from docx import Document
        doc = Document(str(path))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n".join(paragraphs).strip() or None
    except Exception:
        return None


def extract_text_xlsx(path: Path) -> str | None:
    """Extract text from an XLSX file using openpyxl.

    Args:
        path: Path to the XLSX file.

    Returns:
        Extracted text with sheet names as headers, or None on failure.
    """
    try:
        from openpyxl import load_workbook
        wb = load_workbook(str(path), read_only=True, data_only=True)
        parts: list[str] = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows: list[str] = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    rows.append(" | ".join(cells))
            if rows:
                parts.append(f"[{sheet_name}]\n" + "\n".join(rows))
        wb.close()
        return "\n\n".join(parts).strip() or None
    except Exception:
        return None


def extract_text_csv(path: Path) -> str | None:
    """Extract text from a CSV file.

    Args:
        path: Path to the CSV file.

    Returns:
        CSV content as text, or None on failure.
    """
    try:
        import csv
        text_parts: list[str] = []
        with open(path, 'r', encoding='utf-8', errors='replace') as f:
            reader = csv.reader(f)
            for row in reader:
                if any(cell.strip() for cell in row):
                    text_parts.append(" | ".join(cell.strip() for cell in row))
        return "\n".join(text_parts).strip() or None
    except Exception:
        return None


def extract_text_pptx(path: Path) -> str | None:
    """Extract text from a PPTX file using python-pptx.

    Args:
        path: Path to the PPTX file.

    Returns:
        Extracted text with slide numbers as headers, or None on failure.
    """
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        parts: list[str] = []
        for i, slide in enumerate(prs.slides, 1):
            slide_texts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_texts.append(text)
            if slide_texts:
                parts.append(f"[Slide {i}]\n" + "\n".join(slide_texts))
        return "\n\n".join(parts).strip() or None
    except Exception:
        return None


def extract_text_plain(path: Path) -> str | None:
    """Extract text from a plain text file (TXT, MD, RTF).

    Args:
        path: Path to the text file.

    Returns:
        File content as text, or None on failure.
    """
    try:
        content = path.read_text(encoding='utf-8', errors='replace')
        return content.strip() or None
    except Exception:
        return None


_TEXT_EXTRACTORS: dict[str, callable] = {
    '.pdf': extract_text_pdf,
    '.docx': extract_text_docx,
    '.doc': extract_text_docx,
    '.xlsx': extract_text_xlsx,
    '.csv': extract_text_csv,
    '.pptx': extract_text_pptx,
    '.txt': extract_text_plain,
    '.md': extract_text_plain,
    '.rtf': extract_text_plain,
}


def extract_text_from_file(file_path: Path) -> str | None:
    """Extract text content from a document file.

    Routes to the appropriate extractor based on file extension.

    Args:
        file_path: Path to the document file.

    Returns:
        Extracted text, or None if unsupported or extraction failed.
    """
    ext = file_path.suffix.lower()
    extractor = _TEXT_EXTRACTORS.get(ext)
    if extractor:
        return extractor(file_path)
    return None


# -----------------------------------------------------------------------------
# 5. AI ENGINE & EXECUTION
# -----------------------------------------------------------------------------

def normalize_category(raw: str) -> str:
    """Normalize a raw category string: lowercase, replace spaces with underscores, strip unsafe chars.

    Args:
        raw: Raw category string from user input.

    Returns:
        Normalized string (may be empty if input was all special chars).
    """
    normalized = raw.lower().strip().replace(" ", "_")
    safe_chars = [c for c in normalized if c.isalpha() or c.isdigit() or c in ('_', '-')]
    return "".join(safe_chars).strip('_')


def validate_category(raw_category: str | None) -> tuple[str, bool]:
    """Normalize and validate a category name against the allowed list.

    Args:
        raw_category: User-provided category string (may be None or empty).

    Returns:
        A tuple of (normalized_category, was_invalid). was_invalid is True
        if the input fell back to 'uncategorized'.
    """
    if not raw_category or not str(raw_category).strip():
        return 'uncategorized', True
    normalized = normalize_category(str(raw_category))
    if not normalized:
        return 'uncategorized', True
    if normalized in ALLOWED_CATEGORIES:
        return normalized, False
    return 'uncategorized', True


def sanitize_name(raw_name: str) -> str:
    """Convert a raw AI-generated name into a safe snake_case filename stem.

    Args:
        raw_name: Unprocessed name string from the AI response.

    Returns:
        Cleaned, lowercase, underscore-separated name.
    """
    cleaned = raw_name.lower().replace("grid", "").replace("sequence", "")
    cleaned = cleaned.replace(" ", "_")
    safe = "".join([c for c in cleaned if c.isalpha() or c.isdigit() or c in ('_', '-')]).strip('_')
    if len(safe.split('_')) < 3:
        safe = f"{safe}_media_asset"
    return safe


def apply_case_style(name: str, style: str) -> str:
    """Transform a name string to the specified case style.

    Args:
        name: Input name to transform.
        style: One of 'snake_case', 'camelCase', 'kebab-case', 'pascal_case',
               'lowercase', or 'title_case'.

    Returns:
        Name formatted in the requested case style.
    """
    if style == "snake_case":
        return name.lower().replace("-", "_").replace(" ", "_")
    elif style == "camelCase":
        parts = name.replace("-", "_").replace(" ", "_").split("_")
        return parts[0].lower() + "".join(p.capitalize() for p in parts[1:])
    elif style == "kebab-case":
        return name.lower().replace("_", "-").replace(" ", "-")
    elif style == "pascal_case":
        parts = name.replace("-", "_").replace(" ", "_").split("_")
        return "".join(p.capitalize() for p in parts)
    elif style == "lowercase":
        return name.lower().replace("_", "").replace("-", "").replace(" ", "")
    elif style == "title_case":
        parts = name.replace("-", "_").replace(" ", "_").split("_")
        return " ".join(p.capitalize() for p in parts if p)
    else:
        return name


CASE_STYLE_OPTIONS = ["snake_case", "camelCase", "kebab-case", "pascal_case", "lowercase", "title_case"]

CASE_STYLE_LABELS = {
    "snake_case": "snake_case",
    "camelCase": "camelCase",
    "kebab-case": "kebab-case",
    "pascal_case": "PascalCase",
    "lowercase": "lowercase",
    "title_case": "Title Case",
}


def truncate_filename(name: str, max_chars: int) -> str:
    """Truncate a filename stem to max_chars, stripping trailing separators.

    Args:
        name: Filename stem to truncate.
        max_chars: Maximum character count; 0 or negative means no limit.

    Returns:
        Possibly truncated name with trailing underscores/hyphens removed.
    """
    if max_chars <= 0 or len(name) <= max_chars:
        return name
    return name[:max_chars].rstrip("_-")


def _template_date() -> str:
    """Return today's date in ISO format for use in naming templates."""
    return datetime.date.today().isoformat()


def apply_naming_template(template_string: str, asset_data: dict[str, Any]) -> str:
    """Apply a naming template by substituting placeholders from asset data.

    Args:
        template_string: Template with {category}, {topic}, {description}, {date}.
        asset_data: Dictionary containing asset metadata keys.

    Returns:
        Formatted filename stem, or the fallback name if template produces nothing.
    """
    category = asset_data.get('category', 'uncategorized')
    topic = asset_data.get('topic', '') or ''
    description = asset_data.get('description', '') or ''
    fallback = asset_data.get('new_filename', '')

    if not topic and not description:
        # Strip category prefix from fallback if present
        if fallback.lower().startswith(category.lower() + "_"):
            fallback = fallback[len(category) + 1:]
        return fallback

    result = template_string
    result = result.replace("{category}", category)
    result = result.replace("{topic}", topic)
    result = result.replace("{description}", description)
    result = result.replace("{date}", _template_date())

    while "__" in result:
        result = result.replace("__", "_")
    while "--" in result:
        result = result.replace("--", "-")
    result = result.strip("_-")

    if not result or result == template_string:
        return fallback

    return result


def _parse_ai_response(raw_text: str) -> tuple[dict[str, Any] | None, str | None, str | None]:
    """Parse, validate, and extract JSON from an AI model's raw text response.

    Uses Pydantic validation against AssetAnalysisResponse when possible.
    Falls back to raw JSON parsing if the response doesn't match the schema
    (e.g. partial responses from tests or incomplete model outputs).

    Args:
        raw_text: Raw text returned by the AI model.

    Returns:
        A tuple of (parsed_dict, error_type, error_detail). On success,
        error_type and error_detail are None. On failure, parsed_dict is None.
    """
    import re

    clean_res = raw_text.strip()
    if not clean_res:
        return None, 'empty_response', 'Model returned an empty response'

    candidates: list[str] = []

    if clean_res.startswith("```json"):
        candidates.append(clean_res.split("```json")[1].split("```")[0].strip())
    elif clean_res.startswith("```"):
        candidates.append(clean_res.split("```")[1].split("```")[0].strip())

    candidates.append(clean_res)

    json_block = re.search(r"\{[\s\S]*\}", clean_res)
    if json_block and json_block.group(0) not in candidates:
        candidates.append(json_block.group(0))

    for candidate in candidates:
        try:
            raw_dict = json.loads(candidate)
            try:
                validated = AssetAnalysisResponse.model_validate(raw_dict)
                return validated.model_dump(), None, None
            except ValidationError:
                return raw_dict, None, None
        except json.JSONDecodeError:
            continue

    try:
        raw_dict = json.loads(clean_res)
        return raw_dict, None, None
    except json.JSONDecodeError as exc:
        return None, 'json_parse_error', f'JSON decode failed: {exc}'


# -----------------------------------------------------------------------------
# 5b. AI PROVIDERS (Abstract base + implementations)
# -----------------------------------------------------------------------------

VISION_MODEL_PREFIXES = {
    "llava", "bakllava", "qwen2.5vl", "qwen2.5-vl", "qwen2-vl", "qwen3-vl",
    "minicpm", "cogvlm", "moondream",
    "yi-vl", "gemma3", "xclip", "llama3.2-vision", "llama3.2-11b-vision",
    "llama3.2-90b-vision", "pixtral",
}


def _is_vision_model(name: str) -> bool:
    """Check if a model name matches a known vision-capable model prefix.

    Args:
        name: Model name string to check.

    Returns:
        True if the name starts with a recognized vision model prefix.
    """
    name_lower = name.lower().replace(":", "-")
    for prefix in VISION_MODEL_PREFIXES:
        if name_lower.startswith(prefix.lower()):
            return True
    return False


class AIProvider(ABC):
    def __init__(self) -> None:
        """Initialize default model, text model, and API key slots."""
        self._model = ""
        self._api_key = ""
        self.text_model = TEXT_MODEL_NAME

    @abstractmethod
    def analyze(self, base64_img: str, verbose: bool = False,
                prompt_override: str | None = None) -> dict[str, Any]:
        ...

    def analyze_text(self, text_content: str, verbose: bool = False) -> dict[str, Any]:
        """Analyze text content from a document file.

        Default implementation sends the text as a prompt without images.
        Subclasses may override for provider-specific text handling.

        Args:
            text_content: Extracted text from the document.
            verbose: If True, include raw response in error details.

        Returns:
            Result dict with parsed data or error information.
        """
        prompt = f"Document content:\n\n{text_content[:8000]}\n\n---\n\n{get_active_prompt()}"
        return self._analyze_prompt_only(prompt, verbose)

    def _analyze_prompt_only(self, prompt: str, verbose: bool = False) -> dict[str, Any]:
        """Send a text-only prompt to the model (no images).

        Subclasses should override this for provider-specific API calls.

        Args:
            prompt: Full prompt text to send.
            verbose: If True, include raw response in error details.

        Returns:
            Result dict with parsed data or error information.
        """
        result: dict[str, Any] = {'ok': False, 'data': None, 'error': 'unsupported',
                                  'detail': 'Text analysis not supported by this provider', 'raw_response': None}
        return result

    @abstractmethod
    def health_check(self) -> dict[str, Any]:
        ...

    @abstractmethod
    def available_models(self) -> list[str]:
        ...

    @property
    def model(self) -> str:
        """Return the current model name."""
        return self._model

    @model.setter
    def model(self, value: str) -> None:
        """Set the model name."""
        self._model = value

    @property
    def api_key(self) -> str:
        """Return the current API key."""
        return self._api_key

    @api_key.setter
    def api_key(self, value: str) -> None:
        """Set the API key."""
        self._api_key = value

    def _parse_and_validate(self, raw_text: str) -> dict[str, Any]:
        """Parse raw AI text and validate required response keys.

        Args:
            raw_text: Raw text response from the AI model.

        Returns:
            Result dict with keys 'ok', 'data', 'error', 'detail', 'raw_response'.
        """
        result = {'ok': False, 'data': None, 'error': None, 'detail': None, 'raw_response': raw_text}
        parsed, error_type, detail = _parse_ai_response(raw_text)
        if error_type:
            result['error'] = error_type
            result['detail'] = detail
            return result
        if 'new_filename' not in parsed:
            result['error'] = 'missing_keys'
            result['detail'] = "Response JSON is missing required key 'new_filename'"
            return result
        result['ok'] = True
        result['data'] = parsed
        return result


class OllamaProvider(AIProvider):
    def __init__(self) -> None:
        """Initialize Ollama provider with the configured model name."""
        super().__init__()
        self._model = MODEL_NAME
        self._retries = 2

    def analyze(self, base64_img: str, verbose: bool = False,
                prompt_override: str | None = None) -> dict[str, Any]:
        """Send a base64 image to Ollama for AI analysis with retry logic.

        Args:
            base64_img: Base64-encoded JPEG image data.
            verbose: If True, include raw response in error details.
            prompt_override: Optional custom prompt to use instead of get_active_prompt().

        Returns:
            Result dict with parsed data or error information.
        """
        result = {'ok': False, 'data': None, 'error': None, 'detail': None, 'raw_response': None}
        last_exc = None
        prompt = prompt_override or get_active_prompt()
        for attempt in range(self._retries):
            try:
                response = ollama.generate(
                    model=self._model,
                    prompt=prompt,
                    images=[base64_img],
                    keep_alive=MODEL_KEEP_ALIVE,
                    options={"temperature": MODEL_TEMPERATURE, "num_ctx": MODEL_NUM_CTX}
                )
                raw_text = response.get('response', '')
                parsed = self._parse_and_validate(raw_text)
                if parsed['ok'] or attempt == self._retries - 1:
                    return parsed
                last_exc = parsed.get('detail')
            except (ollama.ResponseError, ConnectionError, TimeoutError, OSError) as exc:
                last_exc = exc
                if attempt < self._retries - 1:
                    continue
                result['error'] = 'ollama_error'
                result['detail'] = f'Ollama request failed: {exc}'
                return result
            except Exception as exc:
                result['error'] = 'ollama_error'
                result['detail'] = f'Unexpected AI error: {exc}'
                return result
        if last_exc:
            result['error'] = 'ollama_error'
            result['detail'] = f'Ollama request failed after retry: {last_exc}'
        return result

    def _analyze_prompt_only(self, prompt: str, verbose: bool = False) -> dict[str, Any]:
        """Send a text-only prompt to Ollama (no images) for document analysis.

        Args:
            prompt: Full prompt text to send.
            verbose: If True, include raw response in error details.

        Returns:
            Result dict with parsed data or error information.
        """
        result: dict[str, Any] = {'ok': False, 'data': None, 'error': None, 'detail': None, 'raw_response': None}
        last_exc = None
        model = self.text_model or self._model
        for attempt in range(self._retries):
            try:
                response = ollama.generate(
                    model=model,
                    prompt=prompt,
                    keep_alive=MODEL_KEEP_ALIVE,
                    options={"temperature": MODEL_TEMPERATURE, "num_ctx": MODEL_NUM_CTX}
                )
                raw_text = response.get('response', '')
                parsed = self._parse_and_validate(raw_text)
                if parsed['ok'] or attempt == self._retries - 1:
                    return parsed
                last_exc = parsed.get('detail')
            except (ollama.ResponseError, ConnectionError, TimeoutError, OSError) as exc:
                last_exc = exc
                if attempt < self._retries - 1:
                    continue
                result['error'] = 'ollama_error'
                result['detail'] = f'Ollama request failed: {exc}'
                return result
            except Exception as exc:
                result['error'] = 'ollama_error'
                result['detail'] = f'Unexpected AI error: {exc}'
                return result
        if last_exc:
            result['error'] = 'ollama_error'
            result['detail'] = f'Ollama request failed after retry: {last_exc}'
        return result

    def health_check(self) -> dict[str, Any]:
        """Verify Ollama server is reachable and responsive.

        Returns:
            Dict with 'ok' boolean and 'message' string.
        """
        try:
            ollama.list()
            return {"ok": True, "message": "Ollama is running."}
        except Exception as exc:
            return {"ok": False, "message": f"Ollama not reachable: {exc}"}

    def available_models(self) -> list[str]:
        """List all models available on the Ollama server.

        Returns:
            List of model name strings.
        """
        try:
            tags = ollama.list()
            models = []
            for m in tags.get('models', []):
                if isinstance(m, dict):
                    name = m.get('name', '')
                elif hasattr(m, 'model'):
                    name = m.model
                else:
                    name = str(m)
                if name:
                    models.append(name)
            return models
        except Exception:
            # Only report models that actually exist on the server. Falling back
            # to the config catalog here made every catalog entry look
            # "installed" when the daemon was down (see audit.md §1).
            return []


class GeminiProvider(AIProvider):
    def analyze(self, base64_img: str, verbose: bool = False) -> dict[str, Any]:
        """Analyze an image using the Google Gemini API.

        Args:
            base64_img: Base64-encoded JPEG image data.
            verbose: If True, include raw response in error details.

        Returns:
            Result dict with parsed data or error information.
        """
        result = {'ok': False, 'data': None, 'error': None, 'detail': None, 'raw_response': None}
        if not self._api_key:
            result['error'] = 'api_key_missing'
            result['detail'] = 'Gemini API key not configured.'
            return result
        try:
            model_name = self._model or "gemini-2.0-flash-001"
            url = f"https://generativelanguage.googleapis.com/v1beta/models/{model_name}:generateContent?key={self._api_key}"
            payload = {
                "contents": [{
                    "parts": [
                        {"text": get_active_prompt()},
                        {"inline_data": {"mime_type": "image/jpeg", "data": base64_img}}
                    ]
                }]
            }
            resp = requests.post(url, json=payload, timeout=60)
            resp.raise_for_status()
            data = resp.json()
            candidates = data.get("candidates", [])
            if not candidates:
                result['error'] = 'gemini_empty_response'
                result['detail'] = 'Gemini returned no candidates.'
                return result
            raw_text = ""
            for part in candidates[0].get("content", {}).get("parts", []):
                raw_text += part.get("text", "")
            return self._parse_and_validate(raw_text)
        except requests.exceptions.RequestException as exc:
            result['error'] = 'gemini_api_error'
            result['detail'] = f'Gemini API request failed: {exc}'
            return result
        except Exception as exc:
            result['error'] = 'gemini_api_error'
            result['detail'] = f'Unexpected Gemini error: {exc}'
            return result

    def health_check(self) -> dict[str, Any]:
        """Check if a Gemini API key is configured.

        Returns:
            Dict with 'ok' boolean and 'message' string.
        """
        return {"ok": bool(self._api_key), "message": "API key set" if self._api_key else "No API key configured"}

    def available_models(self) -> list[str]:
        """List models available for the Gemini provider.

        Returns:
            List of model name strings from config.
        """
        return config.get("model", {}).get("providers", {}).get("gemini", {}).get("models", [])


class OpenAIProvider(AIProvider):
    def __init__(self, base_url: str | None = None) -> None:
        """Initialize OpenAI provider with optional base URL override.

        Args:
            base_url: Custom API base URL, or None for the default OpenAI endpoint.
        """
        super().__init__()
        self._base_url = base_url

    def _make_client(self) -> Any:
        """Create and return an OpenAI client instance.

        Returns:
            Configured openai.OpenAI client.
        """
        kwargs = {"api_key": self._api_key}
        if self._base_url:
            kwargs["base_url"] = self._base_url
        return openai.OpenAI(**kwargs)

    def analyze(self, base64_img: str, verbose: bool = False,
                prompt_override: str | None = None) -> dict[str, Any]:
        """Analyze an image using the OpenAI vision API.

        Args:
            base64_img: Base64-encoded JPEG image data.
            verbose: If True, include raw response in error details.
            prompt_override: Optional custom prompt instead of get_active_prompt().

        Returns:
            Result dict with parsed data or error information.
        """
        result = {'ok': False, 'data': None, 'error': None, 'detail': None, 'raw_response': None}
        if not self._api_key:
            result['error'] = 'api_key_missing'
            result['detail'] = 'API key not configured.'
            return result
        try:
            prompt = prompt_override or get_active_prompt()
            client = self._make_client()
            model_name = self._model or "gpt-4o"
            response = client.chat.completions.create(
                model=model_name,
                messages=[{
                    "role": "user",
                    "content": [
                        {"type": "text", "text": prompt},
                        {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{base64_img}"}}
                    ]
                }],
                response_format={
                    "type": "json_schema",
                    "json_schema": {
                        "name": "asset_analysis",
                        "strict": True,
                        "schema": AssetAnalysisResponse.model_json_schema(),
                    },
                },
                max_tokens=1024
            )
            raw_text = response.choices[0].message.content or ""
            return self._parse_and_validate(raw_text)
        except Exception as exc:
            result['error'] = 'openai_api_error'
            result['detail'] = f'OpenAI API request failed: {exc}'
            return result

    def _analyze_prompt_only(self, prompt: str, verbose: bool = False) -> dict[str, Any]:
        """Send a text-only prompt to an OpenAI-compatible endpoint (no images).

        Used for document/audio analysis through cloud providers and the local
        llama.cpp runtime, which expose the same chat-completions surface.

        Args:
            prompt: Full prompt text to send.
            verbose: If True, include raw response in error details.

        Returns:
            Result dict with parsed data or error information.
        """
        result = {'ok': False, 'data': None, 'error': None, 'detail': None, 'raw_response': None}
        if not self._api_key:
            result['error'] = 'api_key_missing'
            result['detail'] = 'API key not configured.'
            return result
        try:
            client = self._make_client()
            model_name = self.text_model or self._model or "gpt-4o"
            response = client.chat.completions.create(
                model=model_name,
                messages=[{
                    "role": "user",
                    "content": [{"type": "text", "text": prompt}],
                }],
                response_format={
                    "type": "json_schema",
                    "json_schema": {
                        "name": "asset_analysis",
                        "strict": True,
                        "schema": AssetAnalysisResponse.model_json_schema(),
                    },
                },
                max_tokens=1024
            )
            raw_text = response.choices[0].message.content or ""
            return self._parse_and_validate(raw_text)
        except Exception as exc:
            result['error'] = 'openai_api_error'
            result['detail'] = f'OpenAI API request failed: {exc}'
            return result

    def health_check(self) -> dict[str, Any]:
        """Check if an OpenAI API key is configured.

        Returns:
            Dict with 'ok' boolean and 'message' string.
        """
        return {"ok": bool(self._api_key), "message": "API key set" if self._api_key else "No API key configured"}

    def available_models(self) -> list[str]:
        """List models available for the OpenAI provider.

        Returns:
            List of model name strings from config.
        """
        return config.get("model", {}).get("providers", {}).get("openai", {}).get("models", [])


class AnthropicProvider(AIProvider):
    def analyze(self, base64_img: str, verbose: bool = False) -> dict[str, Any]:
        """Analyze an image using the Anthropic Claude vision API.

        Args:
            base64_img: Base64-encoded JPEG image data.
            verbose: If True, include raw response in error details.

        Returns:
            Result dict with parsed data or error information.
        """
        result = {'ok': False, 'data': None, 'error': None, 'detail': None, 'raw_response': None}
        if not self._api_key:
            result['error'] = 'api_key_missing'
            result['detail'] = 'API key not configured.'
            return result
        try:
            client = anthropic.Anthropic(api_key=self._api_key)
            model_name = self._model or "claude-3-5-sonnet-20241022"
            response = client.messages.create(
                model=model_name,
                max_tokens=1024,
                messages=[{
                    "role": "user",
                    "content": [
                        {"type": "text", "text": get_active_prompt()},
                        {"type": "image", "source": {"type": "base64", "media_type": "image/jpeg", "data": base64_img}}
                    ]
                }]
            )
            raw_text = response.content[0].text
            return self._parse_and_validate(raw_text)
        except Exception as exc:
            result['error'] = 'anthropic_api_error'
            result['detail'] = f'Anthropic API request failed: {exc}'
            return result

    def health_check(self) -> dict[str, Any]:
        """Check if an Anthropic API key is configured.

        Returns:
            Dict with 'ok' boolean and 'message' string.
        """
        return {"ok": bool(self._api_key), "message": "API key set" if self._api_key else "No API key configured"}

    def available_models(self) -> list[str]:
        """List models available for the Anthropic provider.

        Returns:
            List of model name strings from config.
        """
        return config.get("model", {}).get("providers", {}).get("anthropic", {}).get("models", [])


class GroqProvider(OpenAIProvider):
    def __init__(self) -> None:
        """Initialize Groq provider with its configured base URL."""
        base = config.get("model", {}).get("providers", {}).get("groq", {}).get("base_url", "https://api.groq.com/openai/v1")
        super().__init__(base_url=base)

    def available_models(self) -> list[str]:
        """List models available for the Groq provider.

        Returns:
            List of model name strings from config.
        """
        return config.get("model", {}).get("providers", {}).get("groq", {}).get("models", [])


class OpenRouterProvider(OpenAIProvider):
    def __init__(self) -> None:
        """Initialize OpenRouter provider with its configured base URL."""
        base = config.get("model", {}).get("providers", {}).get("openrouter", {}).get("base_url", "https://openrouter.ai/api/v1")
        super().__init__(base_url=base)

    def available_models(self) -> list[str]:
        """List models available for the OpenRouter provider.

        Returns:
            List of model name strings from config.
        """
        return config.get("model", {}).get("providers", {}).get("openrouter", {}).get("models", [])


LLAMACPP_DEFAULT_URL = "http://localhost:8080"


def _llamacpp_base_url() -> str:
    """Return the OpenAI-compatible base URL for a local llama.cpp server.

    Reads ``LLAMACPP_BASE_URL`` from the environment (or config) so a
    non-default port/host is supported; defaults to ``127.0.0.1:8080``.
    """
    env_url = os.environ.get("LLAMACPP_BASE_URL", "").strip()
    if env_url:
        return env_url.rstrip("/") + "/v1"
    cfg_url = config.get("model", {}).get("llamacpp", {}).get("base_url", "")
    if cfg_url:
        return cfg_url.rstrip("/") + "/v1"
    return LLAMACPP_DEFAULT_URL + "/v1"


LLAMACPP_RUNTIME_PINNED: tuple[str, ...] = (
    "https://github.com/ggml-org/llama.cpp/releases/download/"
    "b10327/llama-b10327-bin-win-cpu-x64.zip",
    "https://github.com/ggml-org/llama.cpp/releases/download/"
    "b10326/llama-b10326-bin-win-cpu-x64.zip",
)


# Published SHA-256 digests (from the GitHub release pages) for every pinned
# runtime zip, so downloads are verifiable even when the GitHub API is down.
LLAMACPP_RUNTIME_PINNED_DIGESTS: dict[str, str] = {
    "https://github.com/ggml-org/llama.cpp/releases/download/"
    "b10327/llama-b10327-bin-win-cpu-x64.zip":
        "c2781932f9af623c9498a12f002f667d2b668f65e0f19b4401e12b5fe9f860c3",
    "https://github.com/ggml-org/llama.cpp/releases/download/"
    "b10326/llama-b10326-bin-win-cpu-x64.zip":
        "dabff645e0948feae41ee6c8e46f2c12dffee96b0cb050e850da7d6b3932f56d",
}


def _llamacpp_runtime_digests() -> dict[str, str]:
    """Return {download_url: sha256-hex} for every llama.cpp runtime candidate.

    The latest release's digest is read from the GitHub API ``digest`` field
    (published as ``sha256:<hex>``). Pinned fallbacks carry the verified
    digests from ``LLAMACPP_RUNTIME_PINNED_DIGESTS`` so setup can verify them
    even when the API is unreachable.

    Returns:
        Map of download URL to lowercase 64-char SHA-256 hex digest.
    """
    digests: dict[str, str] = {}
    try:
        resp = requests.get(
            "https://api.github.com/repos/ggml-org/llama.cpp/releases/latest",
            timeout=10,
        )
        data = resp.json()
        tag = data.get("tag_name", "")
        for asset in data.get("assets", []):
            name = asset.get("name", "")
            if f"{tag}-bin" in name and "-win-cpu-x64" in name and name.endswith(".zip"):
                url = asset.get("browser_download_url", "")
                digest = asset.get("digest", "")
                if url and digest.startswith("sha256:"):
                    digests[url] = digest[len("sha256:"):].strip().lower()
                break
    except Exception:
        pass
    digests.update({k: v.lower() for k, v in LLAMACPP_RUNTIME_PINNED_DIGESTS.items()})
    return digests


def _llamacpp_runtime_urls() -> list[str]:
    """Return candidate download URLs for the llama.cpp Windows CPU runtime.

    Resolves the latest release from GitHub's API and prefers its win-cpu-x64
    zip (the SourceForge-style flexible approach mirrors ExifTool's resolver),
    with pinned release builds as fallbacks so setup never hardcodes a version
    that could go stale.

    Returns:
        List of candidate zip URLs, most-recent first.
    """
    urls: list[str] = []
    try:
        resp = requests.get(
            "https://api.github.com/repos/ggml-org/llama.cpp/releases/latest",
            timeout=10,
        )
        data = resp.json()
        tag = data.get("tag_name", "")
        for asset in data.get("assets", []):
            name = asset.get("name", "")
            if f"{tag}-bin" in name and "-win-cpu-x64" in name and name.endswith(".zip"):
                download_url = asset.get("browser_download_url", "")
                if download_url:
                    urls.append(download_url)
    except Exception:
        pass
    urls.extend(url for url in LLAMACPP_RUNTIME_PINNED if url)
    return urls


class LlamaCppProvider(OpenAIProvider):
    """Runtime provider for a local llama.cpp ``llama-server``.

    Uses the same OpenAI-compatible surface as ``OpenAIProvider`` pointed at
    the local server. A dummy API key satisfies the OpenAI client; llama-server
    does not authenticate. ``available_models()`` lists only the models the
    running server has loaded, mirroring Ollama's honest reporting.
    """

    def __init__(self) -> None:
        """Point an OpenAI-compatible client at the local llama-server."""
        super().__init__(base_url=_llamacpp_base_url())
        self._api_key = "local"
        self._model = MODEL_NAME

    def available_models(self) -> list[str]:
        """List models advertised by the running llama-server.

        Returns:
            List of model name strings, or [] if the server is unreachable.
        """
        try:
            client = self._make_client()
            resp = client.models.list()
            return [m.id for m in getattr(resp, "data", [])]
        except Exception:
            return []


def register_provider(name: str, cls: type[AIProvider]) -> None:
    """Register a provider class in the global provider registry.

    Args:
        name: Short identifier for the provider (e.g. 'ollama').
        cls: Provider class that subclasses AIProvider.
    """
    PROVIDER_REGISTRY[name] = cls


def get_provider(name: str) -> AIProvider:
    """Instantiate and configure a registered provider by name.

    Args:
        name: Provider identifier registered via register_provider.

    Returns:
        Configured AIProvider instance with API key and model set.

    Raises:
        ValueError: If name is not in the provider registry.
    """
    cls = PROVIDER_REGISTRY.get(name)
    if not cls:
        raise ValueError(f"Unknown provider: {name}")
    inst = cls()
    local_runtime = name in ("ollama", "llamacpp")
    if not local_runtime:
        # llama.cpp needs no credentials — keep the provider's placeholder key.
        inst.api_key = load_api_key(name)
    pconf = config.get("model", {}).get("providers", {}).get(name, {})
    valid_models = pconf.get("models", [])
    saved_model = pconf.get("selected_model", "")
    if saved_model and (local_runtime or saved_model in valid_models):
        inst.model = saved_model
    elif not local_runtime and valid_models:
        inst.model = valid_models[0]
    return inst


def list_providers() -> list[str]:
    """Return all registered provider names."""
    return list(PROVIDER_REGISTRY.keys())


register_provider("ollama", OllamaProvider)
register_provider("gemini", GeminiProvider)
register_provider("openai", OpenAIProvider)
register_provider("anthropic", AnthropicProvider)
register_provider("groq", GroqProvider)
register_provider("openrouter", OpenRouterProvider)
register_provider("llamacpp", LlamaCppProvider)


def analyze_asset_with_ai(
    base64_img: str,
    verbose: bool = False,
    retry: bool = True,
    audio_transcription: str | None = None,
) -> dict[str, Any]:
    """Analyze an image using the default Ollama provider.

    Args:
        base64_img: Base64-encoded JPEG image data.
        verbose: If True, include raw response in error details.
        retry: Unused, kept for API compatibility.
        audio_transcription: Optional audio transcription text to include in the prompt.

    Returns:
        Result dict with parsed data or error information.
    """
    provider = get_provider("ollama")
    provider.model = config["model"]["name"]
    prompt_override = None
    if audio_transcription:
        prompt_override = (
            f"Audio transcription (if available):\n{audio_transcription}\n\n---\n\n"
            f"{get_active_prompt()}"
        )
    return provider.analyze(base64_img, verbose=verbose, prompt_override=prompt_override)


def analyze_document_with_ai(text_content: str, verbose: bool = False) -> dict[str, Any]:
    """Analyze document text using the default Ollama provider.

    Args:
        text_content: Extracted text from the document.
        verbose: If True, include raw response in error details.

    Returns:
        Result dict with parsed data or error information.
    """
    provider = get_provider("ollama")
    provider.model = config["model"]["name"]
    provider.text_model = config["model"].get("text_model", TEXT_MODEL_NAME)
    return provider.analyze_text(text_content, verbose=verbose)


def _format_ai_error(ai_result: dict[str, Any], verbose: bool = False) -> str:
    """Format an AI result error into a human-readable message string.

    Args:
        ai_result: Result dict from an AI provider analyze() call.
        verbose: If True, append a snippet of the raw model response.

    Returns:
        Formatted error message string.
    """
    error_type = ai_result.get('error', 'unknown')
    detail = ai_result.get('detail', 'Unknown error')
    messages = {
        'json_parse_error': f'AI response was not valid JSON -- {detail}',
        'missing_keys': detail,
        'empty_response': detail,
        'ollama_error': detail,
        'api_key_missing': detail,
        'gemini_empty_response': detail,
        'gemini_api_error': detail,
        'openai_api_error': detail,
        'anthropic_api_error': detail,
    }
    msg = messages.get(error_type, detail)
    if verbose and ai_result.get('raw_response'):
        snippet = ai_result['raw_response'][:500]
        msg += f"\n    [verbose] Raw model response: {snippet!r}"
    return msg


# -----------------------------------------------------------------------------
# 5c. PER-FORMAT DOCUMENT METADATA
# -----------------------------------------------------------------------------


def _write_docx_metadata(target_file: Path, title: str, summary: str, tags: list[str]) -> None:
    """Write metadata to a DOCX file via python-docx core properties.

    Args:
        target_file: Path to the DOCX file.
        title: Title string to write.
        summary: Description/subject string to write.
        tags: List of keyword strings to write.
    """
    try:
        from docx import Document
        doc = Document(str(target_file))
        doc.core_properties.title = title
        doc.core_properties.subject = summary
        doc.core_properties.keywords = ", ".join(tags)
        doc.save(str(target_file))
    except Exception:
        pass


def _write_xlsx_metadata(target_file: Path, title: str, summary: str, tags: list[str]) -> None:
    """Write metadata to an XLSX file via openpyxl document properties.

    Args:
        target_file: Path to the XLSX file.
        title: Title string to write.
        summary: Description/subject string to write.
        tags: List of keyword strings to write.
    """
    try:
        from openpyxl import load_workbook
        wb = load_workbook(str(target_file))
        wb.properties.title = title
        wb.properties.subject = summary
        wb.properties.keywords = ", ".join(tags)
        wb.save(str(target_file))
    except Exception:
        pass


def _write_document_metadata(target_file: Path, asset: dict[str, Any]) -> bool:
    """Write metadata to a document file using the appropriate format handler.

    Routes to per-format writers (DOCX, XLSX) or skips for formats with
    no standard metadata support (TXT, MD, RTF).

    Args:
        target_file: Path to the committed document file.
        asset: Staged asset dict with tags, summary, staged_name.

    Returns:
        True if metadata was written, False if skipped or failed.
    """
    suffix = target_file.suffix.lower()
    title = asset['staged_name'].replace("_", " ").replace("-", " ").title()

    if suffix in ('.docx', '.doc'):
        _write_docx_metadata(target_file, title, asset['summary'], asset['tags'])
        return True
    if suffix == '.xlsx':
        _write_xlsx_metadata(target_file, title, asset['summary'], asset['tags'])
        return True
    if suffix in ('.txt', '.md', '.rtf', '.csv', '.pptx'):
        return False
    return False


def _build_commit_args(asset: dict[str, Any], target_file: Path) -> list[str]:
    """Build ExifTool argument list for a single asset's metadata.

    Returns an empty list for documents handled by native Python libraries
    (DOCX, XLSX) or formats with no metadata support (TXT, MD, RTF).

    Args:
        asset: Staged asset dict with tags, summary, staged_name.
        target_file: Resolved path to the committed file.

    Returns:
        List of ExifTool CLI arguments (flags + file path as last element),
        or empty list if ExifTool should not be used for this file.
    """
    tag_string = ", ".join(asset['tags'])
    summary = asset['summary']
    title = asset['staged_name'].replace("_", " ").replace("-", " ").title()
    suffix = target_file.suffix.lower()

    if suffix in DOCUMENT_EXTENSIONS and suffix != '.pdf':
        return []

    if suffix in AUDIO_EXTENSIONS:
        if suffix == '.mp3':
            args = [
                "-overwrite_original",
                f"-ID3:TIT2={title}",
                f"-ID3:TALB={summary}",
                f"-ID3:TCOM={summary}",
            ]
            for t in asset['tags']:
                args.append(f"-ID3:TSRC={t}")
            args.append(str(target_file))
            return args
        if suffix in ('.aiff', '.ape'):
            args = [
                "-overwrite_original",
                f"-ID3:TIT2={title}",
                f"-ID3:TALB={summary}",
            ]
            for t in asset['tags']:
                args.append(f"-ID3:TSRC={t}")
            args.append(str(target_file))
            return args
        if suffix in ('.wav', '.flac', '.ogg', '.wv'):
            args = [
                "-overwrite_original",
                f"-XMP-dc:Title={title}",
                f"-XMP-dc:Description={summary}",
            ]
            for t in asset['tags']:
                args.append(f"-XMP-dc:Subject={t}")
            args.append(str(target_file))
            return args
        if suffix in ('.m4a', '.aac'):
            args = [
                "-overwrite_original",
                f"-QuickTime:Title={title}",
                f"-QuickTime:Comment={summary}",
                f"-QuickTime:Keywords={tag_string}",
                str(target_file),
            ]
            return args
        return []

    is_video = suffix in VIDEO_EXTENSIONS

    args = [
        "-overwrite_original",
        "-api", "LargeFileSupport=1",
        f"-XMP-dc:Title={title}",
        f"-XMP-dc:Description={summary}",
        f"-Microsoft:Category={tag_string}",
    ]
    for t in asset['tags']:
        args.append(f"-XMP-dc:Subject={t}")

    if is_video:
        args.extend([
            f"-QuickTime:Title={title}",
            f"-QuickTime:Description={summary}",
            f"-QuickTime:Comment={summary}",
            f"-QuickTime:Keywords={tag_string}",
            f"-Keys:Description={summary}",
            f"-Keys:Keywords={tag_string}",
        ])
    else:
        args.extend([
            f"-EXIF:XPTitle={title}",
            f"-EXIF:XPKeywords={tag_string}",
            f"-Description={summary}",
            f"-Comment={summary}",
        ] + [f"-Keywords={t}" for t in asset['tags']])

    args.append(str(target_file))
    return args


def execute_commit(
    asset: dict[str, Any],
    target_dir: Path,
    sort_into_folders: bool,
    exiftool_session: Any,
    skip_rename: bool = False,
    skip_metadata: bool = False,
) -> str | Path:
    """Rename/move a staged asset to the target directory and write metadata.

    Routes metadata writing by format: ExifTool for images/video/PDF,
    native Python libraries for DOCX/XLSX, skip for TXT/MD/RTF/CSV/PPTX.

    Args:
        asset: Staged asset dict with original_path, staged_name, category, tags, summary.
        target_dir: Destination directory for the file.
        sort_into_folders: If True, create a subfolder named after the category.
        exiftool_session: Active ExifToolSession for writing metadata.
        skip_rename: If True, copy instead of rename (keeps original name).
        skip_metadata: If True, skip all metadata writing.

    Returns:
        Relative path to the committed file, or 'ERROR:<message>' on failure.
    """
    old_path = asset['original_path']
    safe_name = asset['staged_name']
    suffix = old_path.suffix.lower()

    final_folder = target_dir / asset['category'] if sort_into_folders else target_dir
    final_folder.mkdir(parents=True, exist_ok=True)

    new_filename = f"{safe_name}{suffix}"
    new_path = final_folder / new_filename

    counter = 1
    while new_path.exists() and new_path != old_path:
        new_filename = f"{safe_name}_{counter}{suffix}"
        new_path = final_folder / new_filename
        counter += 1

    try:
        if skip_rename:
            final_folder.mkdir(parents=True, exist_ok=True)
            target_file = final_folder / old_path.name
            if old_path != target_file:
                shutil.copy2(str(old_path), str(target_file))
        else:
            old_path.rename(new_path)
            target_file = new_path

        args = _build_commit_args(asset, target_file)
        if not skip_metadata:
            if args:
                exiftool_session.execute(args)
            else:
                _write_document_metadata(target_file, asset)

        if skip_rename:
            return target_file
        return new_path.relative_to(target_dir)
    except Exception as e:
        return f"ERROR:{e}"


def execute_commit_batch(
    assets: list[dict[str, Any]],
    target_dir: Path,
    sort_into_folders: bool,
    exiftool_session: Any,
    skip_rename: bool = False,
    skip_metadata: bool = False,
) -> list[str | Path]:
    """Commit multiple assets in a batch with a single ExifTool IPC call.

    Performs file moves/copies, then sends all metadata writes in one
    -execute block to ExifTool, reducing IPC overhead from ~200ms x N
    to ~200ms total. Documents use native Python metadata writers.

    Args:
        assets: List of staged asset dicts.
        target_dir: Destination directory for file operations.
        sort_into_folders: Whether to sort committed files into category subfolders.
        exiftool_session: Active ExifToolSession for writing metadata.
        skip_rename: If True, copy instead of rename (keeps original name).
        skip_metadata: If True, skip all metadata writing.

    Returns:
        List of relative paths or 'ERROR:<message>' strings, one per asset.
    """
    prepared: list[tuple[dict[str, Any], Path]] = []
    results: list[str | Path] = []

    for asset in assets:
        old_path = asset['original_path']
        safe_name = asset['staged_name']
        suffix = old_path.suffix.lower()
        final_folder = target_dir / asset['category'] if sort_into_folders else target_dir
        final_folder.mkdir(parents=True, exist_ok=True)
        new_filename = f"{safe_name}{suffix}"
        new_path = final_folder / new_filename
        counter = 1
        while new_path.exists() and new_path != old_path:
            new_filename = f"{safe_name}_{counter}{suffix}"
            new_path = final_folder / new_filename
            counter += 1
        try:
            if skip_rename:
                final_folder.mkdir(parents=True, exist_ok=True)
                target_file = final_folder / old_path.name
                if old_path != target_file:
                    shutil.copy2(str(old_path), str(target_file))
            else:
                old_path.rename(new_path)
                target_file = new_path
            prepared.append((asset, target_file))
        except Exception as exc:
            results.append(f"ERROR:{exc}")

    all_args = [_build_commit_args(asset, tf) for asset, tf in prepared]
    if not skip_metadata:
        exiftool_args = [a for a in all_args if a]
        if exiftool_args:
            exiftool_session.execute_batch(exiftool_args)
        for (asset, tf), file_args in zip(prepared, all_args):
            if not file_args:
                _write_document_metadata(tf, asset)

    for asset, target_file in prepared:
        if skip_rename:
            results.append(target_file)
        else:
            final_folder = target_dir / asset['category'] if sort_into_folders else target_dir
            results.append(target_file.relative_to(target_dir))

    return results


# -----------------------------------------------------------------------------
# 5b. SESSION PERSISTENCE
# -----------------------------------------------------------------------------

SESSION_DIR = Path(os.environ.get('APPDATA', Path.home())) / "ai-media-renamer" / "sessions"


def save_session(staged_assets: list[dict[str, Any]], uploaded_files: dict[str, Path], settings: dict[str, Any]) -> Path:
    """Save session state to a JSON file for later restoration."""
    SESSION_DIR.mkdir(parents=True, exist_ok=True)
    ts = datetime.datetime.now().strftime("%Y-%m-%d_%H%M%S")
    session_path = SESSION_DIR / f"session_{ts}.json"

    serializable_assets = []
    for a in staged_assets:
        entry = dict(a)
        entry["original_path"] = str(entry["original_path"])
        serializable_assets.append(entry)

    serializable_files = {name: str(p) for name, p in uploaded_files.items()}

    data = {
        "version": 1,
        "created": ts,
        "staged_assets": serializable_assets,
        "uploaded_files": serializable_files,
        "settings": settings,
    }
    session_path.write_text(json.dumps(data, indent=2, ensure_ascii=False), encoding="utf-8")
    return session_path


def list_sessions() -> list[dict[str, Any]]:
    """Return list of saved sessions sorted by date (newest first)."""
    if not SESSION_DIR.exists():
        return []
    sessions = sorted(SESSION_DIR.glob("session_*.json"), reverse=True)
    result = []
    for s in sessions:
        try:
            data = json.loads(s.read_text(encoding="utf-8"))
            asset_count = len(data.get("staged_assets", []))
            created = data.get("created", s.stem.replace("session_", ""))
            result.append({"path": s, "created": created, "asset_count": asset_count})
        except Exception:
            continue
    return result


def delete_session(session_path: str | Path) -> bool:
    """Delete a saved session file. Returns True on success."""
    try:
        Path(session_path).unlink(missing_ok=True)
        return True
    except Exception:
        return False


def load_session(session_path: str | Path) -> dict[str, Any]:
    """Load a saved session, validating that original files still exist on disk."""
    data = json.loads(Path(session_path).read_text(encoding="utf-8"))

    staged_assets = []
    missing_files = []
    for a in data.get("staged_assets", []):
        a["original_path"] = Path(a["original_path"])
        if a["original_path"].exists():
            staged_assets.append(a)
        else:
            missing_files.append(a["original_name"])

    uploaded_files = {}
    for name, path_str in data.get("uploaded_files", {}).items():
        p = Path(path_str)
        if p.exists():
            uploaded_files[name] = p

    settings = data.get("settings", {})
    return {
        "staged_assets": staged_assets,
        "uploaded_files": uploaded_files,
        "settings": settings,
        "missing_files": missing_files,
    }


# -----------------------------------------------------------------------------
# 5c. UNDO / ROLLBACK
# -----------------------------------------------------------------------------

UNDO_DIR = Path(os.environ.get("APPDATA", Path.home())) / "ai-media-renamer"
UNDO_LOG_FILE = UNDO_DIR / "undo_log.jsonl"


def log_commit_batch(
    batch_id: str,
    target_dir: str,
    records: list[dict[str, Any]],
) -> Path:
    """Write a commit batch to the undo log for potential rollback.

    Args:
        batch_id: Unique identifier for this commit batch (UUID).
        target_dir: Destination directory the files were committed to.
        records: List of per-asset dicts with original_path, new_path, category, tags.

    Returns:
        Path to the undo log file.
    """
    UNDO_DIR.mkdir(parents=True, exist_ok=True)
    entry = {
        "batch_id": batch_id,
        "timestamp": datetime.datetime.now().isoformat(),
        "target_dir": str(target_dir),
        "records": records,
    }
    with open(UNDO_LOG_FILE, "a", encoding="utf-8") as f:
        f.write(json.dumps(entry) + "\n")
    return UNDO_LOG_FILE


def list_undo_batches() -> list[dict[str, Any]]:
    """Return a list of undoable commit batches (most recent first)."""
    if not UNDO_LOG_FILE.exists():
        return []
    batches = []
    with open(UNDO_LOG_FILE, "r", encoding="utf-8") as f:
        for line in f:
            line = line.strip()
            if line:
                try:
                    batches.append(json.loads(line))
                except json.JSONDecodeError:
                    pass
    return list(reversed(batches))


def rollback_last_batch() -> dict[str, Any]:
    """Revert the most recent commit batch.

    Moves files back to original paths and strips injected metadata.

    Returns:
        Dict with ok, restored, failed, errors, batch_id keys.
    """
    batches = list_undo_batches()
    if not batches:
        return {"ok": False, "restored": 0, "failed": 0, "errors": ["No undo batches found"], "batch_id": None}

    batch = batches[0]
    batch_id = batch["batch_id"]
    records = batch.get("records", [])

    restored = 0
    failed = 0
    errors: list[str] = []

    exif = ExifToolSession()

    for rec in records:
        orig = Path(rec["original_path"])
        new = Path(rec["new_path"])

        try:
            if new.exists():
                if new != orig:
                    orig.parent.mkdir(parents=True, exist_ok=True)
                    shutil.move(str(new), str(orig))

                tag_names = rec.get("injected_tags", [])
                if tag_names:
                    tag_args = []
                    for t in tag_names:
                        tag_args.append(f"-{t}=")
                    tag_args.append(str(orig))
                    try:
                        exif.execute(tag_args)
                    except Exception:
                        pass

                restored += 1
            else:
                failed += 1
                errors.append(f"File not found: {new}")
        except Exception as e:
            failed += 1
            errors.append(f"{orig.name}: {e}")

    try:
        exif.close()
    except Exception:
        pass

    remaining = [b for b in batches if b["batch_id"] != batch_id]
    UNDO_DIR.mkdir(parents=True, exist_ok=True)
    with open(UNDO_LOG_FILE, "w", encoding="utf-8") as f:
        for b in reversed(remaining):
            f.write(json.dumps(b) + "\n")

    return {
        "ok": failed == 0,
        "restored": restored,
        "failed": failed,
        "errors": errors,
        "batch_id": batch_id,
    }


# -----------------------------------------------------------------------------
# 5d. DUPLICATE DETECTION
# -----------------------------------------------------------------------------

def compute_asset_hash(file_path: str | Path) -> str | None:
    """Compute hash for an asset. For images/videos, returns perceptual hash (pHash).
    For documents, returns SHA-256 of content. Returns None on failure."""
    path = Path(file_path)
    suffix = path.suffix.lower()

    if suffix in VIDEO_EXTENSIONS:
        try:
            import imagehash as _ih
        except ImportError:
            return None
        h = _compute_video_hash(path, _ih)
        return f"phash:{h}" if h else None

    if suffix in IMAGE_EXTENSIONS:
        try:
            import imagehash as _ih
        except ImportError:
            return None
        h = _compute_image_hash(path, _ih, None)
        return f"phash:{h}" if h else None

    if suffix in AUDIO_EXTENSIONS:
        return _compute_audio_hash(path)

    if suffix in DOCUMENT_EXTENSIONS:
        return _compute_document_hash(path)

    return None


def _compute_document_hash(path: Path) -> str | None:
    """Compute SHA-256 hash for a document file.
    For text-based files under 1MB, hashes extracted text for finer dedup.
    For larger or binary files, hashes raw bytes."""
    try:
        text_suffixes = {'.txt', '.md', '.rtf'}
        if path.suffix.lower() in text_suffixes and path.stat().st_size <= 1_048_576:
            text = extract_text_from_file(path)
            if text:
                return f"sha256:{hashlib.sha256(text.encode('utf-8')).hexdigest()}"
        return f"sha256:{hashlib.sha256(path.read_bytes()).hexdigest()}"
    except Exception:
        return None


def _compute_audio_hash(path: Path) -> str | None:
    """Compute Chromaprint fingerprint for an audio file.

    Uses fpcalc (bundled or system) to generate a perceptual audio fingerprint.
    Returns 'chromaprint:{fingerprint_hex}:{duration}' or None on failure.
    """
    try:
        fpcalc = _resolve_binary_path("fpcalc")
        if not fpcalc:
            return None
        cmd = [fpcalc, "-json", str(path)]
        proc = subprocess.run(cmd, capture_output=True, text=True, timeout=30, creationflags=_NO_WINDOW)
        if proc.returncode != 0:
            return None
        import json as _json
        data = _json.loads(proc.stdout)
        fp = data.get("fingerprint")
        dur = data.get("duration", 0)
        if not fp:
            return None
        return f"chromaprint:{','.join(map(str, fp))}:{dur}"
    except Exception:
        return None


def _compute_image_hash(path: Path, imagehash: Any, pil_image: Any) -> str | None:
    """Hash a single image file."""
    cmd = [
        'ffmpeg', '-y', '-hide_banner', '-loglevel', 'error',
        '-i', str(path),
        '-vf', 'scale=256:256:force_original_aspect_ratio=decrease',
        '-frames:v', '1', '-f', 'image2pipe', '-vcodec', 'mjpeg', '-'
    ]
    try:
        proc = subprocess.run(cmd, capture_output=True, creationflags=_NO_WINDOW)
        if proc.returncode != 0 or not proc.stdout:
            return None
        import io
        from PIL import Image
        img = Image.open(io.BytesIO(proc.stdout))
        return str(imagehash.phash(img))
    except Exception:
        return None


def _compute_video_hash(path: Path, imagehash: Any) -> str | None:
    """Extract midpoint frame from video and compute pHash."""
    duration = get_video_duration(path)
    mid = max(1.0, duration * 0.5)
    cmd = [
        'ffmpeg', '-y', '-hide_banner', '-loglevel', 'error',
        '-ss', str(mid),
        '-i', str(path),
        '-vframes', '1',
        '-vf', 'scale=256:256:force_original_aspect_ratio=decrease',
        '-f', 'image2pipe', '-vcodec', 'mjpeg', '-'
    ]
    try:
        proc = subprocess.run(cmd, capture_output=True, creationflags=_NO_WINDOW)
        if proc.returncode != 0 or not proc.stdout:
            return None
        from PIL import Image
        import io
        img = Image.open(io.BytesIO(proc.stdout))
        return str(imagehash.phash(img))
    except Exception:
        return None


def find_duplicates(staged_assets: list[dict[str, Any]], threshold: int = 10) -> list[dict[str, Any]]:
    """Compare all staged assets pairwise. Returns list of duplicate pairs.
    threshold: Hamming distance threshold for pHash (0-64). Lower = stricter match.
    Default 10 means ~84% similarity required. For SHA-256, exact match only.
    For Chromaprint, uses a similarity threshold of 0.85."""
    try:
        import imagehash as _ih
    except ImportError:
        _ih = None

    raw_hashes: dict[int, str] = {}
    for i, asset in enumerate(staged_assets):
        h = compute_asset_hash(asset["original_path"])
        if h:
            raw_hashes[i] = h

    phash_entries: dict[int, Any] = {}
    sha256_entries: dict[int, str] = {}
    chromaprint_entries: dict[int, str] = {}

    for idx, h in raw_hashes.items():
        if h.startswith("phash:") and _ih is not None:
            try:
                phash_entries[idx] = _ih.hex_to_hash(h[6:])
            except Exception:
                pass
        elif h.startswith("sha256:"):
            sha256_entries[idx] = h[7:]
        elif h.startswith("chromaprint:"):
            chromaprint_entries[idx] = h

    duplicates: list[dict[str, Any]] = []

    phash_indices = sorted(phash_entries.keys())
    for i in range(len(phash_indices)):
        for j in range(i + 1, len(phash_indices)):
            idx_a, idx_b = phash_indices[i], phash_indices[j]
            dist = phash_entries[idx_a] - phash_entries[idx_b]
            if dist <= threshold:
                confidence = max(0, round((1 - dist / 64) * 100))
                duplicates.append({
                    "index_a": idx_a,
                    "index_b": idx_b,
                    "name_a": staged_assets[idx_a]["original_name"],
                    "name_b": staged_assets[idx_b]["original_name"],
                    "hash_type": "phash",
                    "distance": dist,
                    "confidence": confidence,
                })

    sha256_indices = sorted(sha256_entries.keys())
    for i in range(len(sha256_indices)):
        for j in range(i + 1, len(sha256_indices)):
            idx_a, idx_b = sha256_indices[i], sha256_indices[j]
            if sha256_entries[idx_a] == sha256_entries[idx_b]:
                duplicates.append({
                    "index_a": idx_a,
                    "index_b": idx_b,
                    "name_a": staged_assets[idx_a]["original_name"],
                    "name_b": staged_assets[idx_b]["original_name"],
                    "hash_type": "sha256",
                    "distance": 0,
                    "confidence": 100,
                })

    chromaprint_indices = sorted(chromaprint_entries.keys())
    for i in range(len(chromaprint_indices)):
        for j in range(i + 1, len(chromaprint_indices)):
            idx_a, idx_b = chromaprint_indices[i], chromaprint_indices[j]
            sim = _chromaprint_similarity(
                chromaprint_entries[idx_a], chromaprint_entries[idx_b]
            )
            if sim >= 0.85:
                duplicates.append({
                    "index_a": idx_a,
                    "index_b": idx_b,
                    "name_a": staged_assets[idx_a]["original_name"],
                    "name_b": staged_assets[idx_b]["original_name"],
                    "hash_type": "chromaprint",
                    "distance": 0,
                    "confidence": round(sim * 100),
                })

    return duplicates


def _chromaprint_similarity(fp_a: str, fp_b: str) -> float:
    """Compute similarity between two Chromaprint fingerprints.

    Both strings have format 'chromaprint:<ints>:<duration>'.
    Returns 0.0-1.0 similarity score. 0.0 = completely different, 1.0 = identical.
    """
    try:
        _, ints_a, _ = fp_a.split(":", 2)
        _, ints_b, _ = fp_b.split(":", 2)
        a = list(map(int, ints_a.split(",")))
        b = list(map(int, ints_b.split(",")))
        if not a or not b:
            return 0.0
        min_len = min(len(a), len(b))
        if min_len == 0:
            return 0.0
        matches = sum(1 for k in range(min_len) if a[k] == b[k])
        return matches / min_len
    except Exception:
        return 0.0


# -----------------------------------------------------------------------------
# 6. BOOTSTRAP & ENVIRONMENT
# -----------------------------------------------------------------------------


def _resolve_binary_path(name: str) -> str | None:
    """Resolve the filesystem path of a binary, checking PyInstaller bundle first.

    Args:
        name: Binary name (e.g. 'ffmpeg', 'exiftool').

    Returns:
        Absolute path to the binary, or None if not found.
    """
    meipass = getattr(sys, '_MEIPASS', None)
    if meipass:
        candidate = os.path.join(meipass, 'bin', name)
        if os.path.isfile(candidate):
            return candidate
    resolved = shutil.which(name)
    return resolved


def check_ollama_health() -> dict[str, Any]:
    """Probe the Ollama server for connectivity and list vision-capable models.

    Returns:
        Dict with 'connected', 'models', 'all_models', counts, and 'error'.
    """
    try:
        tags = ollama.list()
        models = tags.get('models', [])
        all_names = []
        vision_names = []
        for m in models:
            name = m.get('name', '') if isinstance(m, dict) else str(m)
            if name:
                all_names.append(name)
                if _is_vision_model(name):
                    vision_names.append(name)
        return {
            "connected": True,
            "models": vision_names,
            "all_models": all_names,
            "model_count": len(all_names),
            "vision_count": len(vision_names),
            "error": None,
        }
    except Exception as exc:
        return {
            "connected": False,
            "models": [],
            "all_models": [],
            "model_count": 0,
            "vision_count": 0,
            "error": str(exc),
        }


def _llamacpp_server_running() -> bool:
    """Return True if a local llama.cpp ``llama-server`` answers on its API port."""
    import urllib.parse
    base = _llamacpp_base_url()
    try:
        host_port = base.split("/v1")[0]
        parsed = urllib.parse.urlsplit(host_port)
        host = parsed.hostname or "127.0.0.1"
        port = parsed.port or 8080
        scheme = parsed.scheme or "http"
        url = f"{scheme}://{host}:{port}"
        resp = requests.get(f"{url}/v1/models", timeout=2)
        return resp.status_code == 200
    except Exception:
        return False


def check_environment(profile: list[str] | None = None) -> dict[str, Any]:
    """Verify all required tools and services are available.

    Args:
        profile: Use-case keys (from SETUP_USE_CASES). When provided, only the
            dependencies that profile actually needs are treated as critical —
            e.g. a documents-only user is not blocked by a missing FFmpeg.

    Returns:
        Dict with availability flags for ffmpeg, exiftool, Ollama, and error list.
    """
    profile = profile or []
    needs = use_cases_needs(profile) if profile else {"ffmpeg", "exiftool"}
    ffmpeg_path = _resolve_binary_path("ffmpeg")
    exiftool_path = _resolve_binary_path("exiftool")
    ollama_running = False
    model_available = False
    vision_models: list[str] = []
    text_models: list[str] = []
    errors = []

    if not ffmpeg_path and "ffmpeg" in needs:
        errors.append("FFmpeg not found. Install FFmpeg and add it to your PATH.")

    if not exiftool_path and "exiftool" in needs:
        errors.append("ExifTool not found. Install ExifTool and add it to your PATH.")

    try:
        tags = ollama.list()
        ollama_running = True
        models = tags.get('models', [])
        for m in models:
            if isinstance(m, dict):
                name = m.get('name', '')
            elif hasattr(m, 'model'):
                name = m.model
            else:
                name = str(m)
            if _is_vision_model(name):
                vision_models.append(name)
                model_available = True
            else:
                text_models.append(name)
        # A documents-only profile can run on a text model alone.
        if not model_available and profile and models and "vision_model" not in needs:
            model_available = True
    except Exception:
        ollama_running = False
        errors.append("Ollama is not running. Start Ollama and try again.")

    # Fallback runtime: a local llama.cpp llama-server (OpenAI-compatible).
    llamacpp_running = _llamacpp_server_running()
    if not ollama_running and llamacpp_running:
        provider = get_provider("llamacpp")
        local_models = provider.available_models()
        for name in local_models:
            if _is_vision_model(name):
                vision_models.append(name)
                model_available = True
            else:
                text_models.append(name)
        if not model_available and profile and local_models and "vision_model" not in needs:
            model_available = True

    cloud_configured = CURRENT_PROVIDER not in ("ollama", "llamacpp")

    return {
        "ffmpeg": bool(ffmpeg_path),
        "exiftool": bool(exiftool_path),
        "ollama_running": ollama_running,
        "llamacpp_running": llamacpp_running,
        "model_available": model_available,
        "vision_models": vision_models,
        "text_models": text_models,
        "text_model_available": bool(text_models),
        "cloud_configured": cloud_configured,
        "errors": errors,
    }


def stream_model_download(model_name: str = "qwen2.5vl:7b") -> Any:
    """Stream download progress for an Ollama model pull.

    Args:
        model_name: Ollama model tag to download.

    Yields:
        Dicts with 'status' key and progress/message details.
    """
    try:
        current_stream = ollama.pull(model_name, stream=True)
        for chunk in current_stream:
            status = chunk.get('status', '')
            if status == 'success':
                yield {"status": "success", "message": f"Model {model_name} ready"}
                return

            completed = chunk.get('completed', 0) or 0
            total = chunk.get('total', 0) or 0
            if total and completed:
                percentage = (completed / total) * 100.0
                yield {
                    "status": "progress",
                    "completed": completed,
                    "total": total,
                    "percentage": percentage,
                    "detail": status,
                }
            else:
                yield {"status": "status", "detail": status,
                       "completed": completed, "total": total}

        yield {"status": "success", "message": f"Model {model_name} ready"}
    except Exception as exc:
        yield {"status": "error", "message": str(exc)}


def _parse_version(version: str) -> tuple[int, int, int]:
    """Parse a version string like 'v1.6.0' or '1.4.1' into a comparable tuple."""
    import re
    match = re.match(r"[vV]?(\d+)\.(\d+)\.(\d+)", version.strip())
    if not match:
        return (0, 0, 0)
    return tuple(int(g) for g in match.groups())  # type: ignore[return-value]


def check_for_updates() -> dict[str, Any]:
    """Check GitHub releases for a newer version of the application.

    Returns:
        Dict with 'current', 'latest', 'update_available', 'download_url', 'ok'.
    """
    try:
        resp = requests.get(
            "https://api.github.com/repos/Abdulmusawwir/ai-media-renamer/releases/latest",
            timeout=5
        )
        data = resp.json()
        latest = data.get("tag_name", "")
        return {
            "current": VERSION,
            "latest": latest,
            "update_available": bool(latest) and _parse_version(latest) > _parse_version(VERSION),
            "download_url": data.get("html_url", ""),
            "ok": True,
        }
    except Exception as exc:
        return {"ok": False, "current": VERSION, "latest": "", "update_available": False,
                "download_url": "", "error": str(exc)}


def sha256_file(path: str | Path, chunk_size: int = 1024 * 1024) -> str:
    """Return the lowercase hex SHA-256 digest of a file, streaming in chunks.

    Args:
        path: File to hash (may be large, e.g. multi-GB GGUF models).
        chunk_size: Read chunk size in bytes.

    Returns:
        Lowercase 64-char hex digest.
    """
    hasher = hashlib.sha256()
    with open(path, "rb") as fh:
        for chunk in iter(lambda: fh.read(chunk_size), b""):
            hasher.update(chunk)
    return hasher.hexdigest()


def verify_sha256(path: str | Path, expected: str, label: str) -> None:
    """Raise if a file's SHA-256 digest does not match the published value.

    Args:
        path: File to verify.
        expected: Expected lowercase 64-char hex digest (whitespace-tolerant).
        label: Human-readable name for error messages.

    Raises:
        RuntimeError: When the computed digest differs from ``expected``.
    """
    expected = expected.strip().lower()
    if len(expected) != 64 or any(c not in "0123456789abcdef" for c in expected):
        raise RuntimeError(f"invalid expected SHA-256 digest for {label}: {expected}")
    actual = sha256_file(path)
    if actual != expected:
        raise RuntimeError(
            f"SHA-256 mismatch for {label}: expected {expected}, got {actual}"
        )


def download_file(url: str, dest: Path,
                  progress_callback: Callable[[int, int], None] | None = None,
                  chunk_size: int = 8192,
                  expected_sha256: str | None = None) -> bool:
    """Download a file from an HTTPS URL with optional progress callback.

    HTTPS is mandatory: plain-HTTP downloads are rejected outright. When
    ``expected_sha256`` is supplied, the downloaded bytes are verified against
    the published digest before the file is accepted; a mismatch raises
    ``RuntimeError`` and the partial file is removed.

    Args:
        url: HTTPS URL to download.
        dest: Destination file path.
        progress_callback: Optional function receiving (bytes_downloaded, total_bytes).
        chunk_size: Read chunk size in bytes.
        expected_sha256: Published SHA-256 hex digest to verify against.

    Returns:
        True on successful download.

    Raises:
        ValueError: When ``url`` is not HTTPS.
        RuntimeError: When the downloaded file fails SHA-256 verification.
        requests exceptions on network failure.
    """
    if not url.lower().startswith("https://"):
        raise ValueError(f"Refusing to download over insecure transport: {url}")
    tmp = dest.with_suffix(".part")
    try:
        resp = requests.get(url, stream=True, timeout=30)
        resp.raise_for_status()
        total = int(resp.headers.get("content-length", 0))
        downloaded = 0
        dest.parent.mkdir(parents=True, exist_ok=True)
        with open(tmp, "wb") as f:
            for chunk in resp.iter_content(chunk_size=chunk_size):
                if chunk:
                    f.write(chunk)
                    downloaded += len(chunk)
                    if progress_callback and total:
                        progress_callback(downloaded, total)
        tmp.rename(dest)
        if expected_sha256:
            try:
                verify_sha256(dest, expected_sha256, dest.name)
            except Exception:
                dest.unlink(missing_ok=True)
                raise
        return True
    except Exception:
        if tmp.exists():
            tmp.unlink(missing_ok=True)
        raise


# -----------------------------------------------------------------------------
# 5d. LLAMA.CPP RUNTIME INSTALL + LIFECYCLE (wizard-driven default runtime)
# -----------------------------------------------------------------------------

# GGUF files are stored under the user data dir so they survive app re-installs.
LLAMACPP_MODELS_DIR = USER_DATA_DIR / "models"
LLAMACPP_SERVER_EXE = "llama-server.exe"

LLAMACPP_GGUF_FILENAMES: dict[str, tuple[str, str]] = {
    "qwen2.5vl:7b": ("qwen2-vl-7b-q4_k_m.gguf", "mmproj-qwen2-vl-7b-q8_0.gguf"),
    "qwen2.5vl:2b": ("qwen2-vl-2b-q4_k_m.gguf", "mmproj-qwen2-vl-2b-q8_0.gguf"),
    "qwen2.5:3b": ("qwen2.5-3b-instruct-q4_k_m.gguf", ""),
}


def _llamacpp_gguf_paths(model_name: str) -> tuple[Path, Path]:
    """Return (gguf_path, mmproj_path) where a GGUF catalog model lives.

    Args:
        model_name: Catalog model name (e.g. 'qwen2.5vl:7b').

    Returns:
        (Path to the main GGUF, Path to the mmproj — empty Path when the
        model is text-only).
    """
    gguf_file, mmproj_file = LLAMACPP_GGUF_FILENAMES.get(
        model_name, (model_name.replace(":", "-") + ".gguf", "")
    )
    mmproj = LLAMACPP_MODELS_DIR / mmproj_file if mmproj_file else Path("")
    return LLAMACPP_MODELS_DIR / gguf_file, mmproj


def configure_llamacpp_install(model_name: str, gguf_path: Path,
                               mmproj_path: Path | None = None,
                               make_default: bool = True) -> None:
    """Record a wizard-installed llama.cpp runtime in config and persist it.

    Args:
        model_name: Model name the server advertises (e.g. 'qwen2.5vl:7b').
        gguf_path: Path to the loaded GGUF model file.
        mmproj_path: Optional vision-projector file (mmproj).
        make_default: Also make the app start on the llama.cpp runtime with
            this as its active model.
    """
    config["model"]["llamacpp"] = {
        "base_url": LLAMACPP_DEFAULT_URL,
        "gguf_name": model_name,
        "gguf_path": str(gguf_path),
        "mmproj_path": str(mmproj_path) if mmproj_path else "",
    }
    providers = config["model"].setdefault("providers", {})
    providers.setdefault("llamacpp", {})
    providers["llamacpp"]["models"] = [model_name]
    providers["llamacpp"]["selected_model"] = model_name
    if make_default:
        config["model"]["name"] = model_name
        config["model"]["text_model"] = model_name
        config["model"]["last_provider"] = "llamacpp"
    save_config()


def ensure_llamacpp_server(timeout: int = 40) -> bool:
    """Start the configured llama.cpp server if it is not already running.

    Uses the model paths recorded by ``configure_llamacpp_install``. Returns
    True when a server answers on the API port (either already running or
    freshly started).

    Args:
        timeout: Maximum seconds to wait for the server to come up.
    """
    if _llamacpp_server_running():
        return True
    exe = _resolve_binary_path(LLAMACPP_SERVER_EXE)
    if not exe:
        bundled = USER_DATA_DIR / "bin" / LLAMACPP_SERVER_EXE
        if bundled.exists():
            exe = str(bundled)
    if not exe:
        return False
    llamacpp = config.get("model", {}).get("llamacpp", {})
    gguf = llamacpp.get("gguf_path", "")
    if not gguf or not Path(gguf).exists():
        return False
    args = [exe, "-m", gguf, "--host", "127.0.0.1", "--port", "8080",
            "-c", str(MODEL_NUM_CTX)]
    alias = llamacpp.get("gguf_name", "")
    if alias:
        args += ["--alias", alias]
    mmproj = llamacpp.get("mmproj_path", "")
    if mmproj and Path(mmproj).exists():
        args += ["--mmproj", mmproj]
    try:
        subprocess.Popen(
            args,
            creationflags=_NO_WINDOW | getattr(subprocess, "DETACHED_PROCESS", 0),
            stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL,
        )
        import time as _time
        deadline = _time.time() + timeout
        while _time.time() < deadline:
            if _llamacpp_server_running():
                return True
            _time.sleep(1)
    except Exception:
        return False
    return False


def wait_for_ollama_service(timeout: int = 120) -> bool:
    """Poll the Ollama HTTP endpoint until it responds or timeout is reached.

    Args:
        timeout: Maximum seconds to wait.

    Returns:
        True if Ollama responded within the timeout.
    """
    import time
    url = "http://localhost:11434/api/tags"
    deadline = time.time() + timeout
    while time.time() < deadline:
        try:
            resp = requests.get(url, timeout=2)
            if resp.status_code == 200:
                return True
        except requests.ConnectionError:
            pass
        time.sleep(2)
    return False


def switch_ai_provider(new_provider: str, api_key: str | None = None) -> dict[str, Any]:
    """Switch the active AI provider, release old resources, and persist choice.

    Args:
        new_provider: Provider name to switch to (e.g. 'ollama', 'gemini').
        api_key: Optional API key to store for the new provider.

    Returns:
        Dict with 'ok', 'message', and optionally 'require_download'.
    """
    global CURRENT_PROVIDER, CURRENT_API_KEY

    if CURRENT_PROVIDER == "ollama" and new_provider != "ollama":
        try:
            ollama.generate(model=config["model"]["name"], keep_alive=0)
        except Exception:
            pass

    CURRENT_PROVIDER = new_provider

    provider = get_provider(new_provider)
    if api_key:
        CURRENT_API_KEY = api_key
        save_api_key(new_provider, api_key)
        provider.api_key = api_key
    elif new_provider == "llamacpp":
        # Local llama-server needs no credentials; keep the provider default.
        CURRENT_API_KEY = None
    elif new_provider != "ollama":
        stored = load_api_key(new_provider)
        CURRENT_API_KEY = stored
        provider.api_key = stored
    pconf = config.get("model", {}).get("providers", {}).get(new_provider, {})
    provider.model = pconf.get("selected_model", "") or (pconf.get("models") or [None])[0] or config["model"]["name"]

    config["model"]["last_provider"] = new_provider
    save_config()

    if new_provider not in ("ollama", "llamacpp"):
        return {"ok": True, "message": f"Switched to {new_provider}. Local model weights released from RAM/VRAM."}

    if new_provider == "llamacpp":
        env = check_environment()
        if not env["llamacpp_running"]:
            return {"ok": False, "require_download": False,
                    "message": "llama.cpp server is not running at localhost:8080."}
        if not env["model_available"]:
            return {"ok": False, "require_download": False,
                    "message": "No vision model available. Load a vision GGUF in llama-server."}
        return {"ok": True, "message": "Switched to local llama.cpp."}

    env = check_environment()
    if not env["ollama_running"]:
        return {"ok": False, "require_download": False, "message": "Ollama is not running. Start Ollama first."}
    if not env["model_available"]:
        return {"ok": False, "require_download": True, "message": "No vision model found. Download required."}
    return {"ok": True, "message": "Switched to local Ollama."}


def wipe_local_model(model_name: str = "qwen2.5vl:7b") -> dict[str, Any]:
    """Delete a local Ollama model to free disk space.

    Args:
        model_name: Ollama model tag to remove.

    Returns:
        Dict with 'ok' and 'message'.
    """
    try:
        ollama.delete(model_name)
        return {"ok": True, "message": f"Model {model_name} deleted."}
    except Exception as exc:
        return {"ok": False, "message": str(exc)}


# -----------------------------------------------------------------------------
# 7. STAGING EXPORT / IMPORT
# -----------------------------------------------------------------------------

def export_staging_csv(staged_assets: list[dict[str, Any]]) -> str:
    """Serialize staged assets to a CSV string for export.

    Args:
        staged_assets: List of staged asset dicts.

    Returns:
        CSV-formatted string with all asset fields.
    """
    import csv
    import io
    output = io.StringIO()
    writer = csv.writer(output)
    writer.writerow(["original_name", "proposed_filename", "category", "tags", "summary"])
    for a in staged_assets:
        writer.writerow([
            a.get("original_name", ""),
            a.get("staged_name", ""),
            a.get("category", ""),
            ", ".join(a.get("tags", [])),
            a.get("summary", ""),
        ])
    return output.getvalue()


def export_staging_json(staged_assets: list[dict[str, Any]]) -> str:
    """Serialize staged assets to a JSON string for export.

    Args:
        staged_assets: List of staged asset dicts.

    Returns:
        Pretty-printed JSON string.
    """
    clean = []
    for a in staged_assets:
        clean.append({
            "original_name": a.get("original_name", ""),
            "proposed_filename": a.get("staged_name", ""),
            "category": a.get("category", ""),
            "tags": a.get("tags", []),
            "summary": a.get("summary", ""),
        })
    return json.dumps(clean, indent=2)


def import_staging_csv(csv_string: str, allowed_categories: tuple[str, ...] | list[str]) -> tuple[list[dict[str, Any]], list[str]]:
    """Parse a CSV string into staged asset dicts with category validation.

    Args:
        csv_string: CSV-formatted string with asset rows.
        allowed_categories: Valid category names for validation.

    Returns:
        A tuple of (assets_list, warnings_list) where warnings describe
        any category mismatches that fell back to 'uncategorized'.
    """
    import csv
    import io
    assets = []
    warnings = []
    reader = csv.DictReader(io.StringIO(csv_string))
    allowed = set(allowed_categories)
    for row_num, row in enumerate(reader, start=2):
        tags_raw = row.get("tags", "")
        tags = [t.strip() for t in tags_raw.split(",") if t.strip()]
        category = row.get("category", "").strip().lower().replace(" ", "_")
        safe_chars = [c for c in category if c.isalpha() or c.isdigit() or c in ("_", "-")]
        category = "".join(safe_chars).strip("_")
        if category and category not in allowed:
            warnings.append(f"Row {row_num}: unknown category '{category}' → fallback to 'uncategorized'")
            category = "uncategorized"
        if not category:
            category = "uncategorized"
        assets.append({
            "original_name": row.get("original_name", ""),
            "staged_name": row.get("proposed_filename", ""),
            "category": category,
            "tags": tags,
            "summary": row.get("summary", ""),
        })
    return assets, warnings
